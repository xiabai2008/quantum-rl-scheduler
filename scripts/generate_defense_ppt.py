"""生成答辩PPT文件（Issue #543）。

基于 答辩PPT大纲.md 的18页大纲，使用 python-pptx 生成基础 .pptx 文件。
生成的PPT包含标题、要点和关键数据表格，可作为团队进一步美化的基础。
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ============================================================
# 配色方案（参考大纲附录）
# ============================================================
COLOR_PRIMARY = RGBColor(0x00, 0x66, 0xFF)  # 天衍云蓝
COLOR_SECONDARY = RGBColor(0x8B, 0x5C, 0xF6)  # 量子紫
COLOR_SUCCESS = RGBColor(0x10, 0xB9, 0x81)  # 成功绿
COLOR_WARNING = RGBColor(0xF5, 0x9E, 0x0B)  # 警告橙
COLOR_DARK = RGBColor(0x1E, 0x29, 0x3B)  # 深色文字
COLOR_GRAY = RGBColor(0x64, 0x74, 0x8B)  # 灰色文字
COLOR_LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)  # 浅灰背景
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    """添加封面幻灯片。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    # 标题
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.alignment = PP_ALIGN.CENTER
    # 副标题
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(20)
    p2.font.color.rgb = COLOR_GRAY
    p2.alignment = PP_ALIGN.CENTER


def add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    sub_bullets: list[list[str]] | None = None,
) -> None:
    """添加内容幻灯片（标题+要点列表）。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    # 标题栏背景
    shape = slide.shapes.add_shape(
        1,
        Inches(0),
        Inches(0),
        Inches(10),
        Inches(1.1),  # MSO_SHAPE.RECTANGLE
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()
    # 标题文字
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.LEFT
    # 内容要点
    txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(8.8), Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    sub_bullets = sub_bullets or [[] for _ in bullets]
    for i, (bullet, subs) in enumerate(zip(bullets, sub_bullets, strict=False)):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_DARK
        p.space_after = Pt(8)
        p.level = 0
        for sub in subs:
            sp = tf2.add_paragraph()
            sp.text = sub
            sp.font.size = Pt(16)
            sp.font.color.rgb = COLOR_GRAY
            sp.space_after = Pt(4)
            sp.level = 1


def add_table_slide(
    prs: Presentation,
    title: str,
    intro: str,
    headers: list[str],
    rows: list[list[str]],
    note: str = "",
) -> None:
    """添加带表格的幻灯片。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 标题栏
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    # 介绍文字
    if intro:
        txIntro = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(8.8), Inches(0.6))
        tfIntro = txIntro.text_frame
        pIntro = tfIntro.paragraphs[0]
        pIntro.text = intro
        pIntro.font.size = Pt(16)
        pIntro.font.color.rgb = COLOR_GRAY
    # 表格
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_top = Inches(1.9) if intro else Inches(1.4)
    tbl_height = Inches(0.4 * n_rows)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.8), tbl_top, Inches(8.4), tbl_height
    )
    table = table_shape.table
    # 表头
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = COLOR_WHITE
            paragraph.alignment = PP_ALIGN.CENTER
    # 数据行
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_LIGHT_BG
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = COLOR_DARK
                paragraph.alignment = PP_ALIGN.CENTER
    # 备注
    if note:
        txNote = slide.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(8.8), Inches(0.5))
        tfNote = txNote.text_frame
        tfNote.word_wrap = True
        pNote = tfNote.paragraphs[0]
        pNote.text = note
        pNote.font.size = Pt(11)
        pNote.font.color.rgb = COLOR_GRAY
        pNote.font.italic = True


def generate_ppt(output_path: str) -> None:
    """生成18页答辩PPT。"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # === 第1页：封面 ===
    add_title_slide(
        prs,
        "量子+AI双向赋能：量子RL驱动的智能调度系统",
        '中国电信集团有限公司 2026年度"揭榜挂帅"擂台赛\n选题编号：XA-202609\nAI让量子算力被用好 · 量子让AI更懂真实世界 —— 双向感知闭环',
    )

    # === 第2页：问题定义 ===
    add_content_slide(
        prs,
        "核心矛盾：量子算力稀缺且昂贵 × 调度粗糙（FCFS）造成浪费",
        [
            "问题：量子资源稀缺且昂贵，当前量子云平台普遍采用FCFS调度",
            "  - 导致量子通道拥堵与经典通道空转并存、高负载时低优先级任务饥饿",
            "影响：高优先级量子任务常因排队错过最佳执行窗口（如金融盘前风控）",
            "切入：将强化学习（RL）引入量子云任务调度",
            "  - 实现量子/经典/混合任务智能分流与动态优化",
            "  - 赛题'双向赋能'闭环：AI调度量子资源 + 量子真机数据反哺AI环境校准",
        ],
    )

    # === 第3页：方案概览 ===
    add_content_slide(
        prs,
        "双向感知闭环：AI赋能调度 + 量子真机数据反哺AI",
        [
            "核心思路：AI智能调度量子资源（强证据），量子真机数据驱动AI环境校准（方法论）",
            "第一层：调度层 — PPO智能调度实时最优分流，+20.2% vs 真实FCFS（N=250, p=7.56e-12）",
            "第二层：编译层 — PPO替代SABRE比特映射，深电路SWAP减少+38.5%(N=80, Wilcoxon p=2.75e-02，事后子集方向性提示)",
            "第三层：量子赋能AI — 真机噪声模型提取→仿真环境校准→PPO噪声敏感性评估（N=25配对p=2.98e-08，负向敏感性证据）",
            "双向闭环：AI优化量子调度→量子真机测量数据反哺AI评估环境校准→对真实噪声有感知、可校准的调度决策",
            "关键数字：PPO +20.2%（vs 真实FCFS）、等待时间-14.0%双优、推理延迟~1ms、315次真机调用100%成功",
        ],
    )

    # === 第4页：RL调度引擎（建模）===
    add_content_slide(
        prs,
        "RL调度引擎：16维状态空间 + 毫秒级决策",
        [
            "状态空间（16维）：量子比特可用率、队列长度、平均等待时间、保真度、串扰风险、到达率时序等",
            "动作空间（4类）：经典执行、量子执行、混合执行、量子误差缓解(QEM，动作位预留)",
            "奖励函数设计：多目标奖励",
            "  - 兼容性约束（错配惩罚 -2.0）",
            "  - 执行收益（量子最高 10×speedup）",
            "  - 等待惩罚（超阈值累加）",
            "  - 利用率惩罚（量子闲置 -1.0）",
            "决策延迟：单次前向推理 ~1ms（0.88-1.02ms实测，100-10000任务规模），满足实时调度需求",
        ],
    )

    # === 第5页：8策略对比 ===
    add_table_slide(
        prs,
        "PPO策略综合奖励领先：比真实FCFS提升+20.2%（N=250）",
        "8策略公平对比（N=250，50 seed × 5 episode，16维交付模型）",
        ["策略", "平均奖励", "vs FCFS", "备注"],
        [
            ["PPO", "1982.69", "+20.2%", "第一名, p=7.56e-12"],
            ["FCFS", "1648.91", "基线", "真实 FCFS（量子路由）"],
            ["SJF", "774.86", "-53.0%", "短任务优先"],
            ["DQN(Random占位)", "602.37", "-63.5%", "DQN已删除，Random替代"],
            ["Greedy", "80.71", "-95.1%", "量子优先，均衡场景崩溃"],
            ["Quantum-Only", "-826.59", "-150.1%", "单一资源基线"],
        ],
        "多目标权衡（诚实披露）：PPO 综合奖励 +20.2%（vs 真实 FCFS），量子利用率 -3.3%（未达团队基线 R-P-01 的 30% 目标，不作卖点；该目标非比赛方案明文要求）；统计协议：Welch t + Mann-Whitney U 双验证，28组Bonferroni校正（α=0.0018），95% CI [+14.3%, +26.7%]，功效分析达标，原始JSON可复现",
    )

    # === 第6页：高负载防饥饿 ===
    add_table_slide(
        prs,
        "高负载公平调度：PPO在极端拥塞下的表现",
        "场景：任务到达率 > 吞吐量（λ=1.2），极端拥塞场景；N=5 seeds",
        ["策略", "总奖励", "饥饿数(5seeds)", "Jain资源公平"],
        [
            ["PPO", "2874.25", "69", "0.838"],
            ["FCFS", "1998.00", "74", "0.795"],
            ["SJF", "1035.10", "90", "0.926"],
        ],
        "核心发现：PPO在高负载下总奖励最高（2874.25 vs FCFS 1998.00，+43.9%），饥饿控制与FCFS相当（69 vs 74）；SJF公平性最优但总奖励最低。2026-08-08 以真实FCFS基线（按队首任务类型调度）重测，旧'0.60 vs 22.00'基于弱基线已废弃",
    )

    # === 第7页：编译层优化 ===
    add_content_slide(
        prs,
        "AI赋能编译层：PPO替代SABRE比特映射（公平对比v2）",
        [
            "问题：传统量子电路编译使用启发式算法（SABRE）进行比特映射，效率较低",
            "方案：使用PPO强化学习替代传统启发式算法，直接学习最优比特映射策略",
            "核心数据：公平对比v2（4×4 2D网格拓扑，同池配对60电路，Issue #451）",
            "  - 深电路(14-16q) SWAP减少+38.5%（N=80, Wilcoxon p=2.75e-02显著；全60电路p=8.40e-01不显著，诚实披露）",
            "  - 原76.4%为不公平对比已废弃",
            "  - 整体p=8.40e-01不显著，端到端编译对比（布局+路由）",
            "模型信息：ppo_compilation_agent.zip，200k steps，4×4 2D网格全电路分布训练",
            "叙事：AI从调度到编译的全栈赋能量子计算",
        ],
    )

    # === 第8页：量子赋能AI（主方向）===
    add_content_slide(
        prs,
        "量子赋能AI：真机数据驱动的噪声敏感性评估",
        [
            "核心逻辑：用量子硬件的真实测量数据校准AI评估环境（方法论贡献）",
            "闭环流程：",
            "  - 真机H门测量（保真度0.976）→ 噪声分布建模",
            "  - → 注入仿真环境 → 配对检验评估噪声对AI决策评估的影响",
            "权威证据（N=25配对检验）：",
            "  - 噪声致奖励下降12.43%（p=2.98e-08，d_z=7.71，CI[-19.49%,-4.69%]，事后功效1.0）",
            "  - 方向为负——实证'量子硬件测量结果影响AI决策评估'这一机制成立",
            "探索性方向（noise_feedback_v2，N=50）：训练注入奖励+1.5%（p=0.584不显著，方向性证据）",
            "科学定位：这是'敏感性/机制'证据，不是'量子提升AI性能'——我们如实呈现负向结果，因为实验是真的",
        ],
    )

    # === 第9页：QUBO退火（已降级）===
    add_content_slide(
        prs,
        "探索方向与诚实边界",
        [
            "QUBO退火（探索性，默认关闭）：",
            "  - 将PPO决策头构造为QUBO矩阵，D-Wave neal模拟退火求解",
            "  - 结果：奖励变化-5.6%，p=0.9430（20seeds不显著），训练开销+74.5%",
            "  - 已标记@deprecated，不作为性能claim",
            "量子赋能AI主方向：真机噪声敏感性评估（N=25配对，p=2.98e-08，负向敏感性证据）",
            "其他诚实边界（主动披露）：",
            "  - 量子利用率-3.3%（未达团队基线R-P-01的30%目标，不作卖点）",
            "  - 性能结论以仿真N=250支撑，真机为可用性验证+小样本探索",
            "我们的立场：负结果如实报告，因为实验是真的",
        ],
    )

    # === 第10页：推理延迟验证 ===
    add_table_slide(
        prs,
        "推理延迟验证：PPO决策延迟稳定在~1ms",
        "不同任务规模下的推理延迟对比",
        ["任务规模", "PPO延迟(ms)", "FCFS延迟(ms)", "SJF延迟(ms)"],
        [
            ["100", "0.882", "0.431", "0.372"],
            ["500", "0.904", "0.386", "0.381"],
            ["1000", "1.018", "0.359", "0.368"],
            ["5000", "0.914", "0.373", "0.419"],
            ["10000", "0.897", "0.370", "0.387"],
        ],
        "核心发现：PPO单步决策延迟稳定在~1ms，远低于100ms实时调度要求；O(1)复杂度，不随任务规模增长",
    )

    # === 第11页：真机验证 ===
    add_table_slide(
        prs,
        "真机可用性验证：315次SDK调用100%成功",
        "天衍平台超导量子计算机（目标：天衍-287，105数据比特；历史小样本数据来自 tianyan176 回退）",
        ["策略", "平均奖励", "备注"],
        [
            ["PPO", "1736.32", "N=10 小样本（历史，tianyan176；d=5.33 探索性）"],
            ["SJF", "575.33", "N=10 小样本"],
            ["FCFS", "383.00", "N=10 小样本基线"],
        ],
        "诚实声明：真机验证主要证明SDK可用性（315次100%成功），性能数字以仿真N=250结果为准；上表为小样本探索性数据",
    )

    # === 第12页：多场景压力测试 ===
    add_content_slide(
        prs,
        "多场景压力测试：PPO跨场景适应性最强",
        [
            "4种压力场景：均衡负载、高负载、量子资源波动、混合潮汐",
            "PPO综合稳定性最强：4场景中2次第一、2次第二，平均排名1.8",
            "量子资源波动场景优势明显（方向性观察；历史探索数据基于诚实化前基线，权威 stress 数据待重跑核定）",
            "Greedy两极分化：高负载场景第一，但量子波动场景暴跌至倒数第一",
            "场景-算法适配决策树：",
            "  - 量子资源波动 → PPO（最优）",
            "  - 均衡负载 → PPO（最优）",
            "  - 高负载 → Greedy（略优）或PPO（稳定次优）",
            "  - 混合潮汐 → SJF（最稳定）或PPO（稳定次优）",
        ],
    )

    # === 第13页：系统演示 ===
    add_content_slide(
        prs,
        "Web监控面板：实时调度可视化",
        [
            "技术栈：FastAPI后端 + Vue3前端 + Echarts图表，单页应用无刷新",
            "核心功能：",
            "  - 实时任务队列看板（任务ID、类型、优先级、等待时间）",
            "  - RL Agent决策过程可视化（当前状态→策略输出→选中动作）",
            "  - 资源利用率仪表盘（量子/经典实时跳动）",
            "  - 高负载拥塞分流可视化（λ=1.2场景）",
            "  - 多机器状态监控（3台真机在线状态、负载、最近提交）",
            "演示亮点：任务到达后~1ms内完成决策，面板实时刷新无延迟",
        ],
    )

    # === 第14页：创新点总结 ===
    add_content_slide(
        prs,
        "三大创新点",
        [
            "创新1：AI赋能量子计算调度（核心claim）",
            "  - PPO强化学习实现毫秒级动态决策，+20.2% vs FCFS（p=7.56e-12）",
            "  - 三层架构：调度层+编译层+高负载公平调度（真实FCFS基线重测）",
            "创新2：量子赋能AI（主方向：真机噪声反馈）",
            "  - 真机噪声→模型校准→PPO噪声敏感性评估（负向证据，噪声致奖励下降12.43%），形成量子→AI感知闭环",
            "  - 奖励 +1.5%（p=0.584 不显著）、派单率 p=0.020 显著但效应量极小——诚实定位为方向性证据，不宣称统计成立",
            "创新3：天衍云平台全链路对接",
            "  - 315次真机调用100%成功，三态熔断器+三级降级策略",
        ],
    )

    # === 第15页：工程实现 + 可信度主张 ===
    add_content_slide(
        prs,
        "工程化与可信度：可复现、可验证、有门禁",
        [
            "工程稳定性：三级降级策略、三态熔断器、Prometheus 7项监控",
            "可复现性：requirements.lock 锁定已验证组合（sb3 2.9.0/torch 2.8.0），全新环境一键安装",
            "测试保障：主套件3725用例（全量3746含benchmark）0失败，ruff/mypy/bandit全绿，CI 9项全绿",
            "数字门禁：统计口径一致性检查（含PDF/PPTX二进制扫描）+ 权威数字校验，防止数字漂移",
            "统计协议（严谨性主张）：N=250独立运行、Welch t + MWU双验证、28组Bonferroni校正、",
            "  95% CI + 效应量 + 功效分析、预注册方案、负结果全披露——每个数字可复现",
        ],
    )

    # === 第16页：行业场景（金融）===
    add_content_slide(
        prs,
        "落地场景：量子金融风险计算的调度痛点",
        [
            "场景：券商风险计量部门——盘前风控窗口内运行量子风险计算（蒙特卡洛/组合优化）",
            "痛点：量子算力稀缺、任务异构（量子/经典/混合）、风控时效敏感、机时成本高",
            "系统对症：",
            "  - PPO 1ms决策 → 任务实时分流到量子/经典/混合",
            "  - 优先级加权 → 盘前高价值任务优先",
            "  - 等待惩罚 → 风控窗口内完成任务",
            "  - 公平调度 → 多租户资源均衡",
            "仿真支撑：综合收益+20.2%、等待时间-14.0%（N=250）；真机链路315次100%成功",
            "诚实边界：场景价值为估算（详见 industry_case_finance_v1.md），需真实部署验证；",
            "  我们的价值主张=调度效率（已验证）+ 真机接入能力（已验证），不夸大机时节省",
        ],
    )

    # === 第17页：团队介绍 ===
    add_content_slide(
        prs,
        "团队分工",
        [
            "团队规模：8人",
            "算法组（3人）：RL环境设计、PPO/DQN训练、噪声校准模块",
            "工程组（3人）：天衍云API对接、Mock环境、Web监控面板、CI/CD",
            "产品组（2人）：实验报告、PPT制作、演示视频、文档整理",
            "[待填充]：成员姓名、学校/单位、个人照片",
            "[待填充]：指导老师姓名、单位",
        ],
    )

    # === 第18页：致谢 + Q&A ===
    add_content_slide(
        prs,
        "致谢 & 预设问答",
        [
            "致谢：",
            "  - 感谢中国电信集团有限公司提供揭榜挂帅擂台赛平台",
            "  - 感谢天衍云平台的真机支持（天衍-287超导量子计算机）",
            "  - 感谢指导老师的悉心指导",
            "预设Q&A：",
            "  - Q: 量子→AI是负向结果，双向赋能还成立吗？ A: 双向=双向感知协同，非双向都提升。AI→量子交出+20.2%统计显著成果；量子→AI交出真机数据驱动的方法论——实证'量子硬件测量影响AI评估'机制成立（N=25, p=2.98e-08）。发现'噪声是挑战'本身就是研究探索成果；敢如实报告负结果，正结果才更可信",
            "  - Q: 利用率-3.3%怎么解释？ A: 优化目标是综合调度收益（多目标权衡），单指标优化会牺牲其他维度；等待时间-14.0%双优；如实披露正是因为实验是真的",
            "  - Q: 为什么性能验证靠仿真？ A: 真机验证两级：可用性（315次100%成功）已达成；性能级受免费机时限制为N=10小样本探索。权威性能结论由N=250仿真+真实FCFS基线支撑，完全可复现",
            "  - Q: 有没有真实金融客户？ A: 诚实回答：没有。做的是仿真验证+真机可用性验证，真实部署是下一步。赛题要求'研究与应用探索'，我们的协议和真机链路已为落地打好基础",
        ],
    )

    prs.save(output_path)
    print(f"PPT已生成: {output_path}（共{len(prs.slides)}页）")


if __name__ == "__main__":
    generate_ppt("deliverable_models/答辩PPT.pptx")
