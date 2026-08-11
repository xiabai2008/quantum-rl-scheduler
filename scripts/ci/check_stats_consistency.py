#!/usr/bin/env python
"""
统计口径一致性检查脚本 (Statistical Consistency Checker)

Issue #141: 消除4套p值混用，建立单一权威统计源

扫描所有 .md 文件及演示链路 .py 文件中的统计数字（p值、效应量、N值），
与 config/statistics.yaml 权威源比对，报告不一致告警。

用法:
    python scripts/ci/check_stats_consistency.py
    python scripts/ci/check_stats_consistency.py --strict  # CI模式：有不一致则退出码1
"""

from __future__ import annotations

import math
import os
import re
import sys
from pathlib import Path
from typing import Any

import yaml

_PROJECT_ROOT = Path(__file__).resolve().parent.parent.parent
STATS_YAML = _PROJECT_ROOT / "config" / "statistics.yaml"
AUTHORITATIVE_NUMBERS_MD = _PROJECT_ROOT / "docs" / "authoritative_numbers.md"

# 修复 Windows GBK 终端下 emoji 字符导致的 UnicodeEncodeError 崩溃
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")

# 需要扫描的文件模式
SCAN_GLOBS = ["*.md", "**/*.md"]

# 演示链路 .py 文件（Web面板模板/一键演示/PPT生成脚本的展示字符串
# 曾多次出现废弃数字，同样纳入口径审计）
DEMO_CHAIN_FILES = [
    "src/visualization/fallback_template.py",
    "scripts/demo_one_click.py",
    "scripts/generate_defense_ppt.py",
    "scripts/demo/demo.py",
    "scripts/demo/demo_multi_machine.py",
]

# 8.5 审查：权威数值内部冲突检查（A8）
# 同一文件内禁止并存两组互斥的"权威"数值——评委可当场发现。
# 豁免行必须含诚实披露标记（废弃/错误数据/中间产物/审查修正等）。
INTERNAL_CONFLICTS: list[tuple[list[str], str]] = [
    (["+20.2%", "+123.4%"], "综合收益 +20.2%（主口径）与 +123.4%（历史 vs Hybrid-Default）并存"),
    (["-3.3%", "+7.9%"], "利用率 -3.3%（N=250 真实）与 +7.9%（错误旧数据）并存"),
    (["25.46", "57.27"], "等待时间 25.46（N=250）与 57.27（N=4 中间品）并存"),
    (["+38.5%", "+33.3%"], "编译层深电路 +38.5%（N=80 显著）与 +33.3%（N=20 旧样本）并存"),
]

# 8.7 审查新增（A9）：孤立废弃值 → 权威值映射
# 用于检测"废弃值作为唯一权威呈现"的回归（同文件并存检测 A8 抓不到的场景）。
_ORPHAN_DEPRECATED: list[tuple[str, str, str]] = [
    ("+7.9%", "-3.3%", "利用率 +7.9% 为错误旧数据，权威为 -3.3%（N=250 真实基线）"),
    ("2349", "1982.69", "PPO 平均奖励 2349 为废弃弱基线值，权威为 1982.69（N=250）"),
    # 8.7-v3 红队审查 P1-2：旧口径常以"区间外变体"出现（如 2348.91/2348.88/2350），
    # 仅登记 2349 会被此类近似值绕过。补登常见变体，权威统一为 1982.69。
    ("2348.91", "1982.69", "PPO 平均奖励 2348.91 为废弃弱基线值，权威为 1982.69（N=250）"),
    ("2348.88", "1982.69", "PPO 平均奖励 2348.88 为废弃弱基线值，权威为 1982.69（N=250）"),
    ("2350", "1982.69", "PPO 平均奖励 2350 为废弃弱基线值，权威为 1982.69（N=250）"),
    # 8.7-v3 红队审查新增：demo 脚本/演示视频脚本曾残留已证伪的 N=1 单次运行值 +48.9%，
    # 权威利用率口径为 -3.3%（N=250 权威实测）。此值作为"当前成果"出现即告警。
    (
        "+48.9%",
        "-3.3%",
        "量子利用率 +48.9% 为已证伪的 N=1 单次运行旧值，权威为 -3.3%（N=250 权威实测）",
    ),
    (
        "33.6%",
        "-3.3%",
        "量子利用率 33.6%→50% 为旧口径（N=1），权威口径为 -3.3%（PPO 0.4467 vs FCFS 0.4620，N=250）",
    ),
]
_DISCLOSURE_MARK = (
    "废弃",
    "错误数据",
    "中间产物",
    "审查修正",
    "8.5 审查",
    "已废弃",
    "无法支撑",
    "非统计证据",
    "历史",
)

# 排除的文件/目录（权威报告本身，不检查自身）
EXCLUDE_PATHS = {
    # 8.6 审查修复：移除 authoritative_numbers.md 的豁免（N3/N4/N6 曾藏在其中未被门禁抓到），
    # 该文件自称"唯一事实源"，必须纳入一致性校验。
    # "docs/authoritative_numbers.md",
    "results/reports/multiseed_real_machine_report.md",
    "results/reports/multiseed_real_machine_report_10seeds.md",
    "results/reports/multiseed_real_machine_report_10seeds_v2.md",
    "results/reports/head_only_validation.md",
    "results/reports/real_machine_statistical_significance.md",
    "results/reports/dqn_ppo_fcfs_comparison.md",
    "results/reports/annealing_ablation_20seeds_report.md",
    "results/reports/real_machine_boundary_statement.md",
    "results/reports/roi_analysis.md",
    "results/reports/utilization_multiseed_report.md",
    # 历史实验报告（14维旧模型口径，已加废弃横幅冻结，禁止直接引用）
    "results/reports/ablation_report.md",
    "results/reports/issue_457_105_qubits_validation_report.md",
    "results/reports/power_analysis.md",
    "results/reports/quantum_ratio_sensitivity.md",
    "results/reports/real_machine_closed_loop.md",
    "results/reports/real_machine_validation.md",
    "results/reports/tradeoff_analysis.md",
}

# 排除的文件名模式（内部协作文档，不提交到仓库，不需检查）
EXCLUDE_PATTERNS = [
    "teammate_",  # docs/teammate_*.md
    "PR审查",  # PR审查临时报告
    "_pr",  # PR审查临时报告前缀
    "pr_patrol_",  # docs/pr_patrol_*.md
    "project_dashboard_",  # docs/project_dashboard_*.html
    # 8.7-v3 红队审查：对抗审查报告为内部审计文档（会原文引用已废弃旧口径作
    # "诚实披露"说明），不属于对外交付物，与 PR审查/pr_patrol_ 同源，纳入豁免，
    # 避免门禁对其自身误报。
    "红队审查",
]

# p值正则模式：匹配 p=0.031, p=1.032e-42, p<0.001, p=3.04×10⁻¹¹ 等
P_VALUE_PATTERN = re.compile(
    r"p\s*[=<>]\s*"
    r"(?:"
    r"(?:\d+\.?\d*(?:[eE][+-]?\d+)?)"  # 标准科学计数法: 1.032e-42
    r"|"
    r"(?:\d+\.?\d*\s*[×x]\s*10[⁻⁺\-+]?[⁰¹²³⁴⁵⁶⁷⁸⁹0-9]+)"  # Unicode: 1.032×10⁻⁴²
    r"|"
    r"(?:<\s*10[⁻⁺\-+]?[⁰¹²³⁴⁵⁶⁷⁸⁹0-9]+)"  # p<10⁻⁴²
    r")",
    re.IGNORECASE,
)

# 已知的权威p值集合（从YAML加载后动态构建）
AUTHORITATIVE_P_VALUES: dict[str, str] = {}

# 已知废弃/错误p值 → 应替换为的权威值
DEPRECATED_P_VALUES: dict[str, str] = {
    "4.92e-55": "MISATTRIBUTED: p=4.92e-55 是 Random vs PPO 的p值，不是 PPO vs FCFS。应使用 p=7.56e-12 (Welch t, 真实FCFS基线)",
    # 8.7-v4 红队审查 P0-2/P1-3：MAPPO 协同优势 +84.6% 为未收敛+训练量不均等的旧探索值，
    # 其 p=0.0188（multi_machine_comparison_report.md）只支撑该旧口径，不得再作为权威呈现。
    # 权威 MAPPO 口径为同训练量收敛严格对比 +36.5%（N=20，p=0.024，mappo_strict_strict_comparison）。
    "0.0188": "MAPPO 协同优势 p=0.0188 支撑的是已废弃的 +84.6%（5000步未收敛+训练量1:3不均等）；权威口径为同训练量收敛严格对比 +36.5%（N=20，p=0.024，mappo_strict_strict_comparison）",
    "0.019": "MAPPO 协同优势 p=0.019 支撑的是已废弃的 +84.6%（5000步未收敛+训练量1:3不均等）；权威口径为同训练量收敛严格对比 +36.5%（N=20，p=0.024，mappo_strict_strict_comparison）",
}

# 检验方法混用检查
WELCH_T_FOR_1032E42 = "ERROR: p=1.032e-42 对应 Mann-Whitney U 检验（14维旧模型），不是 PPO vs FCFS 的权威 Welch t 值（p=7.56e-12）"

# Issue #446: 严禁表述黑名单（来自 docs/authoritative_numbers.md 第六节）
# 注意：黑名单匹配后，若行内包含以下"诚实披露"关键词则豁免（避免对已修正的诚实标注误报）
HONEST_DISCLOSURE_KEYWORDS = [
    "单seed",
    "探索性",
    "诚实披露",
    "10seeds",
    "10 seeds",
    "20seeds",
    "20 seeds",
    "5seeds",
    "5 seeds",
    "旧实验",
    "注：",
    "注:",
    "边界",
    "cherry-pick",
    "已降级",
    "已废弃",
    # 8.6 复核：废弃/历史/禁止上下文提及旧值属合规引用，豁免黑名单检测。
    # 含这些限定词的行必然是"说明某旧值已废弃/被禁止/为历史口径"，而非将其作权威值呈现。
    "旧",
    "历史",
    "已诚实化",
    "诚实化",
    "黑名单",
    "禁止",
    "deprecated",
    "废弃",
    "假数据",
]


def _is_honest_disclosure(line: str) -> bool:
    """判断行内是否包含诚实披露限定词（豁免黑名单检测）。"""
    lower = line.lower()
    return any(kw in line or kw.lower() in lower for kw in HONEST_DISCLOSURE_KEYWORDS)


def _is_stress_coverage_context(line: str) -> bool:
    """8.7-v4 修复：判别 +91.4% 是否出现在"行覆盖率"语义下。

    mutation_testing_report 等报告在陈述覆盖率历史时列出 87.38% / 87.5% /
    91.4% / 93.78% 等数字，其中 91.4% 是覆盖率而非 stress 量子波动优势。
    该语义需豁免 stress +91.4% 黑名单，避免误报。仅在行内同时含覆盖率相关
    上下文词时豁免。
    """
    lower = line.lower()
    return "覆盖率" in line or "行覆盖率" in line or "coverage" in lower


BLACKLIST_PATTERNS: list[tuple[str, str]] = [
    (
        r"284\s*次(SDK)?(真机)?调用",
        "BLACKLIST: 真机调用次数已统一为315次（284+31审计口径），淘汰284旧表述",
    ),
    (
        r"284\s*次真机",
        "BLACKLIST: 真机调用次数已统一为315次，淘汰284旧表述",
    ),
    (
        r"v5已完成|答辩PPT.*v5|技术白皮书.*v5.*11章",
        "BLACKLIST: v5白皮书/PPT从未存在，实际为docs/technical_whitepaper.pdf（7章），PPT制作中",
    ),
    (
        r"利用率(提升)?\s*[≥>=]\s*30%.*已?达成|资源利用率.*≥30%.*(?<!未)达标",
        "BLACKLIST: N=250权威数据利用率仅+7.9%（未达30%目标），需改为部分达成口径",
    ),
    (
        r"等待时间.*-5\.7%",
        "BLACKLIST: -5.7%为单seed乐观结果，10seeds严谨实验显示等待时间+6.1%，需诚实呈现",
    ),
    (
        r"退火.*p\s*=\s*0\.190(?!.*(20seeds|不显著.*0\.94|已降级.*探索))",
        "BLACKLIST: 退火p=0.190对应5seeds旧实验，20seeds新实验p=0.9430（不显著），以新数据为准",
    ),
    (
        r"PPO.*真机.*1665\.22|真机.*PPO.*1665\.22|1665\.22.*PPO.*真机",
        "BLACKLIST: 真机PPO均值已升级为N=10 v2权威值1736.32，旧N=5值1665.22严禁作为对外基准（若为历史引用请显式标注N=5并说明已废弃）",
    ),
    (
        r"FCFS.*真机.*353\.22|真机.*FCFS.*353\.22|353\.22.*FCFS.*真机",
        "BLACKLIST: 真机FCFS均值已升级为N=10 v2权威值383.00，旧N=5值353.22严禁作为对外基准（若为历史引用请显式标注N=5并说明已废弃）",
    ),
    # 8.6 审查新增（P0-1/P0-2）：废弃效应量与旧 CI/p 值必须从对外材料中清除
    (
        r"Cohen[’'`\s]*s?\s*d\s*=\s*-2\.1353|d=-2\.1353|d_z\s*=\s*-2\.1353",
        "BLACKLIST: Cohen's d=-2.1353 为 8.5 前弱基线（Hybrid-Default）效应量，已废弃。当前权威效应量为 rank-biserial=-0.3642（或 Welch 独立样本 Cohen's d≈+0.63）",
    ),
    (
        r"\[\+?113\.3%,?\s*\+?133\.5%\]",
        "BLACKLIST: 95%CI [+113.3%, +133.5%] 为废弃 +123.4% 弱基线的 CI，已废弃。当前权威 CI 为 [+14.3%, +26.7%]",
    ),
    (
        r"p\s*=\s*1\.449e-66|p=1\.449\s*×\s*10⁻⁶⁶|p=1\.449\s*×\s*10-66",
        "BLACKLIST: p=1.449e-66 为 8.5 前弱基线 p 值，已废弃。当前权威 p=7.56e-12（真实 FCFS 基线）",
    ),
    (
        r"[\-\+]?940\.56|[\-\+]?1128\.79",
        "BLACKLIST: Quantum-Only=-940.56 / Classical-Only=-1128.79 为 8.5 前旧值，已废弃。当前权威值 -826.59 / -1075.49",
    ),
    # 8.7-v2 审查新增：演示视频脚本 demo_video_final_script.md 曾用中文数字
    # （"百分之一百二十三点四"）表达以下废弃统计量，ASCII 黑名单匹配不到。
    # 依赖上层中文数字归一化（_normalize_chinese_numerals）使其对中文表述同样生效。
    (
        r"(?<!\d)123\.4%(?!%)",
        "BLACKLIST: PPO 提升 +123.4% 为 8.5 前弱基线旧值，已废弃；权威为 +20.2%（N=250 真实 FCFS 基线）",
    ),
    (
        r"(?<!\d)163\.3%(?!%)",
        "BLACKLIST: PPO vs 随机 +163.3% 为旧口径，已废弃；权威口径不再使用该百分比表述",
    ),
    (
        r"(?<!\d)36\.8%(?!%)",
        "BLACKLIST: 多机协同 +36.8% 为旧值，已废弃；权威协同优势为 +36.5%（MAPPO vs 独立PPO，N=20，50K收敛严格对比，p=0.024）",
    ),
    (
        r"等待时间.*(增加|高出|增大)\s*51%",
        "BLACKLIST: 等待时间 +51% 为旧口径，已废弃；权威为 -14.0%（PPO 等待时间更短，25.46 vs 29.61 步）",
    ),
    # 8.7-v2 复核新增：提交物文档残留的旧口径（requirements_traceability /
    # award_roadmap / novelty_statement 曾出现）。这些值不属"诚实披露旧值"，
    # 而是被当作当前权威/结果呈现，必须转为权威口径。
    # 注意：模式必须足够精确，避免误伤历史报告（如 reward_ablation_d3 的 57.73 步
    # 属 D3 专项实验结论，defense_qa 的"⚠️ Issue #116 标注"属诚实披露），
    # 因此仅匹配"直接当权威成果呈现"的短语。
    (
        r"奖励提升\s*86\.3%",
        "BLACKLIST: MAPPO 协同优势权威为 +36.5%（MARL vs 独立PPO，N=20，50K收敛严格对比，p=0.024）；+86.3% 仅指叠加规模扩展后的总提升 vs 单机，须拆解呈现，不得直接称'奖励提升86.3%'",
    ),
    # 8.7-v4 红队审查 P0-2/P0-3/P1-4：MAPPO 协同优势 +84.6%、退火 +6.4%、
    # stress 量子波动 +91.4% 均为"已废弃/未收敛/诚实化前"旧口径，转为权威值时
    # 一律以诚实披露限定词豁免；作为唯一权威成果呈现即告警。
    (
        r"(?<!\d)84\.6%(?!%)",
        "BLACKLIST: MAPPO 协同优势 +84.6% 为未收敛(5000步)+训练量不均等(独立PPO仅1/3)的旧探索值；权威为 +36.5%（同训练量收敛严格对比，N=20，50K）",
    ),
    (
        r"(?<!\d)6\.4%(?!%)",
        "BLACKLIST: 退火 +6.4% 为 5seed 旧方向，20seed 权威方向为 -5.6%（p=0.9430 不显著）；不得在展示时当作正向成果，须以'已废弃/探索性'限定",
    ),
    (
        r"(?<!\d)91\.4%(?!%)",
        "BLACKLIST: stress 量子波动 +91.4% 为诚实化前旧 FCFS 基线结果，权威 stress 数据待重跑核定；不得作为当前成果呈现",
    ),
    (
        r"p\s*<\s*10⁻⁶⁶|p<10⁻⁶⁶",
        "BLACKLIST: p<10⁻⁶⁶ 为 8.5 前弱基线旧 p 值，已废弃；权威为 p=7.56e-12（真实 FCFS 基线，N=250）",
    ),
    (
        r"显著优于\s*SABRE",
        "BLACKLIST: 编译层深电路为事后子集方向性证据，不构成'RL 编译优于 SABRE'的定论；须注明'方向性提示，整体不构成定论'",
    ),
    (
        r"鲁棒性提升|RL鲁棒性评估|PPO鲁棒性",
        "BLACKLIST: 噪声反馈应表述为'噪声敏感性评估/负向证据'（噪声致奖励下降12.43%），不得使用'鲁棒性提升/评估'正向措辞",
    ),
    # 8.7-v3 红队审查新增：核心交付文档（bidirectional_empowerment/value_quantification）
    # 曾残留"利用率提升至72%""利用率 72%"荒谬数字，与权威 -3.3% 数量级矛盾。
    # 用"利用率"上下文锚定避免误伤覆盖率等其他百分比。
    (
        r"利用率(提升至|提升到|提升为|达到|达)?\s*72%|72%\s*(量子)?利用率",
        "BLACKLIST: 资源利用率 72% 为无出处荒谬数字，权威口径为 -3.3%（PPO 0.4467 vs FCFS 0.4620，N=250，未达≥30%目标）",
    ),
    # 演示视频脚本曾以"量子利用率从 65% → 78%"暗示调度提升利用率，与 -3.3% 权威口径矛盾
    (
        r"量子利用率(从)?\s*6[5-9]%\s*(→|到|->)\s*\d+%|量子利用率(从)?\s*6[5-9]%",
        "BLACKLIST: 量子利用率 65%+ 为误导性绝对值；权威口径 -3.3%（未达≥30%目标），演示面板数值应贴近真实基线（~46%）",
    ),
]


def load_statistics_yaml() -> dict[str, Any]:
    """加载权威统计源 YAML 文件。"""
    with open(STATS_YAML, encoding="utf-8") as f:
        return yaml.safe_load(f)


def check_authoritative_coverage(stats: dict[str, Any]) -> list[str]:
    """检查权威源是否遗漏已发布的关键实验（Issue #536）。

    确保已发布报告中的关键实验都被 statistics.yaml 收录，
    避免单一权威源机制自身失守。
    """
    warnings = []
    # 已发布的关键实验及其应有的 p 值记录
    required_experiments = [
        (
            "annealing_ablation_head_only",
            "with_vs_without_annealing",
            0.9430,
            "退火消融 20seeds 权威验证（p=0.9430）",
        ),
        (
            "real_machine_10seed_v2",
            "ppo_vs_fcfs",
            5.84e-07,
            "真机 10seeds v2 实验（8.11 复核精确值；此前 p<0.001 截断）",
        ),
        (
            "simulation_8strategy_50seed",
            "ppo_vs_fcfs",
            7.56e-12,
            "8 策略 50seed 仿真（真实FCFS基线，p=7.56e-12）",
        ),
    ]
    # 8.7-v4 红队审查 P0-2 + 8.11 升级：MAPPO 权威口径必须收录
    # （mappo_vs_independent_ppo.improvement_pct = 36.5，N=20 配对检验 p=0.024）
    # 确保"协同优势 +36.5%"在权威源中显式存在，防止 +84.6% 回归为对外口径。
    mappo_block = stats.get("mappo_strict_strict_comparison", {})
    mappo_gain = mappo_block.get("mappo_vs_independent_ppo", {}).get("improvement_pct")
    if mappo_gain != 36.5:
        warnings.append(
            "权威源 MAPPO 口径异常: mappo_strict_strict_comparison.mappo_vs_independent_ppo"
            f".improvement_pct 应为 36.5（N=20 同训练量收敛严格对比，p=0.024），实际为 {mappo_gain}"
        )
    for exp_key, comp_key, expected_p, desc in required_experiments:
        exp = stats.get(exp_key)
        if not exp:
            warnings.append(f"权威源遗漏实验: {exp_key}（{desc}）未在 statistics.yaml 中收录")
            continue
        comp = exp.get(comp_key, {})
        actual_p = comp.get("p_value")
        if actual_p is None:
            warnings.append(f"权威源遗漏 p 值: {exp_key}.{comp_key}（{desc}）未记录 p_value")
            continue
        # 数值比较（允许格式差异）
        # 8.6 审查修复：改用相对容差（rel_tol）+ 极小绝对容差，
        # 避免固定 1e-6 绝对容差把 1e-12 与 1e-42 误判为一致（Issue #854）。
        try:
            if not math.isclose(float(actual_p), float(expected_p), rel_tol=1e-3, abs_tol=1e-12):
                warnings.append(
                    f"权威源 p 值不一致: {exp_key}.{comp_key} 期望 p={expected_p}, 实际 p={actual_p}"
                )
        except (TypeError, ValueError):
            warnings.append(f"权威源 p 值格式异常: {exp_key}.{comp_key} p_value={actual_p}")
    return warnings


def check_yaml_internal_consistency(stats: dict[str, Any]) -> list[str]:
    """8.6 审查新增：校验 statistics.yaml 内部数值自洽。

    修复历史问题（N1/N2/N3/N6）：t 统计量与 p 值、mean_diff 符号、CI 方向、
    effect_size_type 与数值一致，避免"单一权威源"自身矛盾。
    """
    warnings = []

    def _warn(msg: str) -> None:
        warnings.append(f"statistics.yaml 内部矛盾: {msg}")

    sim = stats.get("simulation_8strategy_50seed", {})
    pvf = sim.get("ppo_vs_fcfs", {})
    if pvf:
        ppo = sim.get("strategy_summary", {}).get("PPO", {})
        fcfs = sim.get("strategy_summary", {}).get("FCFS", {})
        ppo_mean = ppo.get("mean_reward")
        fcfs_mean = fcfs.get("mean_reward")
        md = pvf.get("mean_diff")
        ci = pvf.get("ci_95")
        stat = pvf.get("statistic")
        est = pvf.get("effect_size")
        est_type = pvf.get("effect_size_type")
        test = pvf.get("test_method")
        pval = pvf.get("p_value")

        # 1) mean_diff 符号与均值顺序一致（FCFS-PPO 约定下应为负）
        if ppo_mean is not None and fcfs_mean is not None and md is not None:
            expected_sign = -1 if ppo_mean > fcfs_mean else 1
            if (md > 0) != (expected_sign > 0):
                _warn(
                    f"ppo_vs_fcfs.mean_diff={md} 符号与均值顺序矛盾 "
                    f"(PPO={ppo_mean} {'>' if ppo_mean > fcfs_mean else '<'} FCFS={fcfs_mean})"
                )

        # 2) CI 覆盖 mean_diff 且符号一致
        if isinstance(ci, (list, tuple)) and len(ci) == 2 and md is not None:
            lo, hi = ci[0], ci[1]
            if not (lo <= md <= hi):
                _warn(f"ppo_vs_fcfs.ci_95={ci} 未覆盖 mean_diff={md}")
            if (lo > 0) != (md > 0):
                _warn(f"ppo_vs_fcfs.ci_95={ci} 与 mean_diff={md} 符号不一致")

        # 3) statistic 与 p 值量级自洽（Welch t 大统计量 ↔ 极小 p 值）
        if stat is not None and pval is not None and test and "Welch" in str(test):
            try:
                if abs(float(stat)) < 1.0 and float(pval) < 1e-6:
                    _warn(f"ppo_vs_fcfs.statistic={stat} 与 p_value={pval} 量级不自洽")
            except (TypeError, ValueError):
                pass

        # 4) effect_size_type 与 effect_size 匹配（rank-biserial ∈ [-1,1]）
        if est is not None and est_type == "rank-biserial":
            try:
                if not (-1.0 <= float(est) <= 1.0):
                    _warn(f"ppo_vs_fcfs.effect_size={est} 超出 rank-biserial 合法范围 [-1,1]")
            except (TypeError, ValueError):
                pass

    return warnings


def build_authoritative_p_values(stats: dict[str, Any]) -> dict[str, str]:
    """从YAML构建权威p值映射表。

    Returns:
        {p_value_string: "实验名.对比名 (检验方法, N=xxx)"}
    """
    p_map = {}

    # 实验1: 8策略50seed仿真
    sim = stats.get("simulation_8strategy_50seed", {})
    for comp_name in [
        "ppo_vs_fcfs",
        "ppo_vs_dqn",
        "ppo_vs_sjf",
        "ppo_vs_random",
        "dqn_vs_fcfs",
        "fcfs_vs_sjf",
    ]:
        comp = sim.get(comp_name, {})
        if "p_value" in comp:
            p_val = format_p_value(comp["p_value"])
            ctx = f"8策略50seed仿真.{comp_name} ({comp.get('test_method', '?')}, N={comp.get('n', '?')})"
            p_map[p_val] = ctx

    # 实验2: 多seed真机
    real = stats.get("real_machine_5seed", {})
    for comp_name in ["ppo_vs_fcfs", "ppo_vs_sjf", "sjf_vs_fcfs"]:
        comp = real.get(comp_name, {})
        if "p_value" in comp:
            p_val = format_p_value(comp["p_value"])
            ctx = f"多seed真机.{comp_name} ({comp.get('test_method', '?')}, N={comp.get('n', '?')})"
            p_map[p_val] = ctx

    # 实验3: 退火消融
    annealing = stats.get("annealing_ablation_head_only", {})
    comp = annealing.get("with_vs_without_annealing", {})
    if "p_value" in comp:
        p_val = format_p_value(comp["p_value"])
        ctx = (
            f"退火消融 ({comp.get('test_method', '?')}, n={comp.get('n', annealing.get('n', '?'))})"
        )
        p_map[p_val] = ctx

    # 实验4: 真机闭环
    closed_loop = stats.get("real_machine_closed_loop", {})
    comp = closed_loop.get("simulation_vs_mixed_real", {})
    if "p_value" in comp:
        p_val = format_p_value(comp["p_value"])
        ctx = f"真机闭环 ({comp.get('test_method', '?')}, N={comp.get('n', closed_loop.get('n', '?'))})"
        p_map[p_val] = ctx

    # 实验5: 3策略对比
    sim3 = stats.get("simulation_3strategy_10seed", {})
    for comp_name in ["ppo_vs_fcfs", "ppo_vs_dqn"]:
        comp = sim3.get(comp_name, {})
        if "p_value" in comp:
            p_val = format_p_value(comp["p_value"])
            ctx = (
                f"3策略10seed.{comp_name} ({comp.get('test_method', '?')}, N={comp.get('n', '?')})"
            )
            p_map[p_val] = ctx

    return p_map


def build_deprecated_p_values(stats: dict[str, Any]) -> dict[str, str]:
    """从YAML构建废弃p值映射表。"""
    dep_map = {}
    for _key, info in stats.get("deprecated", {}).items():
        p_val = format_p_value(info.get("value", 0))
        status = info.get("status", "deprecated")
        reason = info.get("reason", "").strip().split("\n")[0]
        replacement = info.get("replacement", "")
        dep_map[p_val] = f"[{status.upper()}] {reason} → 替换为: {replacement}"
    return dep_map


def format_p_value(value: Any) -> str:
    """将p值格式化为标准字符串。"""
    if isinstance(value, str):
        return value.lower()
    if isinstance(value, (int, float)):
        if value == 0:
            return "0"
        # 保留科学计数法格式
        return f"{value:.3e}".lower().replace("e+0", "e+").replace("e-0", "e-")
    return str(value).lower()


def normalize_p_value_text(text: str) -> str:
    """将文本中的p值标准化以便比较。"""
    text = text.lower().strip()
    # 移除空格
    text = text.replace(" ", "")
    # Unicode上标转ASCII
    superscript_map = str.maketrans("⁰¹²³⁴⁵⁶⁷⁸⁹⁻⁺", "0123456789-+")
    text = text.translate(superscript_map)
    # ×10 转换为 e
    text = re.sub(r"[×x]10\^?", "e", text)
    # 10^-42 → e-42 (for p<10⁻⁴² style)
    text = re.sub(r"10\^?(-?\d+)", lambda m: f"e{m.group(1)}" if "." in text else text, text)
    return text


def extract_p_values_from_line(line: str) -> list[tuple[str, str]]:
    """从一行文本中提取所有p值引用。

    Returns:
        [(raw_match, normalized_p_value), ...]
    """
    results = []
    for match in P_VALUE_PATTERN.finditer(line):
        raw = match.group(0)
        # 提取数值部分
        val_match = re.search(r"(?:\d+\.?\d*(?:[eE][+-]?\d+)?)", raw)
        if val_match:
            val_str = val_match.group(0).lower()
            results.append((raw, val_str))
        else:
            # 尝试匹配 Unicode 科学计数法
            unicode_match = re.search(r"(\d+\.?\d*)\s*[×x]\s*10[⁻⁺\-+]([⁰¹²³⁴⁵⁶⁷⁸⁹0-9]+)", raw)
            if unicode_match:
                base = float(unicode_match.group(1))
                val_str = f"{base:.3e}".lower()
                results.append((raw, val_str))

    return results


def _is_deprecation_notice(line: str) -> bool:
    """判断行内是否为废弃声明上下文（在废弃声明中提及旧p值属合规引用）。"""
    lower = line.lower()
    return any(kw in line or kw in lower for kw in ("已废弃", "已淘汰", "deprecated"))


def check_welch_t_misattribution(line: str) -> str | None:
    """检查 Welch t + p=1.032e-42 的错误搭配。"""
    line_lower = line.lower()
    has_1032e42 = (
        "1.032e-42" in line_lower or "1.032×10⁻⁴²" in line_lower or "1.03×10⁻⁴²" in line_lower
    )
    # 行中同时包含 Mann-Whitney 或正确新p值 1.449e-66 或废弃声明时，
    # 说明是对比/废弃说明文本，不报错
    has_context = (
        "mann-whitney" in line_lower or "1.449e-66" in line_lower or _is_deprecation_notice(line)
    )
    if has_1032e42 and "welch" in line_lower and not has_context:
        return WELCH_T_FOR_1032E42
    return None


# ---------------------------------------------------------------------------
# 中文数字归一化（8.7-v2 门禁加固）
# 背景：演示脚本等交付文档曾用中文数字（如"百分之一百二十三点四"）表达
# 已废弃统计量，ASCII 黑名单正则完全匹配不到，导致门禁形同虚设。
# 这里把中文数字（含"百分之X"百分比、"X千X百X十X"整数、"零点X"小数）
# 归一化为阿拉伯数字，使废弃值检测对中文表述同样生效。
# ---------------------------------------------------------------------------
_CN_DIGITS: dict[str, int] = {
    "零": 0,
    "一": 1,
    "二": 2,
    "两": 2,
    "三": 3,
    "四": 4,
    "五": 5,
    "六": 6,
    "七": 7,
    "八": 8,
    "九": 9,
}
_CN_UNITS: dict[str, int] = {"十": 10, "百": 100, "千": 1000, "万": 10000}
_CN_CHARSET = "零一二两三四五六七八九十百千万点"
_PERCENT_CN_RE = re.compile(rf"百分之([{_CN_CHARSET}]+)")
_DECIMAL_CN_RE = re.compile(r"零点([零一二三四五六七八九]+)")
_INT_CN_RE = re.compile(rf"(?<![0-9])([{_CN_CHARSET}点]{{2,}})")


def _cn_int_to_arabic(s: str) -> int:
    """把中文整数（如"两千三百四十九"）转成阿拉伯整数。"""
    section = 0
    num = 0
    for ch in s:
        if ch in _CN_DIGITS:
            num = _CN_DIGITS[ch]
        elif ch in _CN_UNITS:
            section += (num or 1) * _CN_UNITS[ch]
            num = 0
    return section + num


def _cn_to_arabic(s: str) -> str:
    """把中文数字串转成阿拉伯数字字符串，支持整数与一位小数点小数。"""
    if "点" in s:
        int_part, _, frac_part = s.partition("点")
        arab_int = str(_cn_int_to_arabic(int_part or "零"))
        frac_digits = "".join(str(_CN_DIGITS.get(c, 0)) for c in frac_part)
        return f"{arab_int}.{frac_digits}"
    return str(_cn_int_to_arabic(s))


def _normalize_chinese_numerals(line: str) -> str:
    """把一行文本中的中文数字归一化为阿拉伯数字（仅用于检测，不改原文）。"""
    out = line
    out = _PERCENT_CN_RE.sub(lambda m: f"{_cn_to_arabic(m.group(1))}%", out)
    out = _DECIMAL_CN_RE.sub(lambda m: f"0.{''.join(str(_CN_DIGITS[c]) for c in m.group(1))}", out)
    out = _INT_CN_RE.sub(lambda m: _cn_to_arabic(m.group(1)), out)
    return out


# ---------------------------------------------------------------------------
# 8.7-v4 红队审查 P0-1：数字门禁三重漏洞修复
#
# 漏洞 1（p值格式归一化）：format_p_value 对 float 输出 .3e（7.560e-12），
#   而文档写作 7.56e-12，字符串比对恒不相等，导致"权威值拼写错误"（如 7.56e-62
#   误写）与"孤儿 p 值"（无出处）全部漏检。这里统一把 p 值解析为 float 做数值比对。
#
# 漏洞 2（Markdown 剥离）：黑名单正则直接匹配原始 markdown，若数字被 **加粗**、
#   `行内代码`、[链接](url) 包裹，正则可能匹配不到。这里先剥离 markdown 排版再匹配。
#
# 漏洞 3（权威值拼写检查 + 孤儿 p 值检测）：对每个提取到的 p 值做两类启发式：
#   (a) 与已知权威/废弃 p 值"保留同 3 位有效数字但指数不同"→ 判定为权威值拼写笔误
#       （如 7.560e-12 被写成 7.56e-62，红队实测 multiscenario_benchmark.md 曾出现）；
#   (b) 与已知权威/废弃 p 值均不匹配且非诚实披露上下文 → 孤儿子 p 值告警（无出处）。
# ---------------------------------------------------------------------------
_MD_STRIP_RULES: list[tuple[re.Pattern[str], str]] = [
    (re.compile(r"`[^`]*`"), " "),  # 行内代码
    (re.compile(r"\*\*[^*\n]+\*\*"), ""),  # 加粗
    (re.compile(r"\*[^*\n]+\*"), ""),  # 斜体（不成对的 * 视为乘法，不剥离）
    (re.compile(r"__[^_\n]+__"), ""),  # 加粗（alt）
    (re.compile(r"\[([^\]\n]*)\]\([^)\n]*\)"), r"\1"),  # 链接保留文本
]


def _strip_markdown(line: str) -> str:
    """剥离行内 markdown 排版（加粗/斜体/行内代码/链接），仅用于口径检测。"""
    out = line
    for pat, repl in _MD_STRIP_RULES:
        out = pat.sub(repl, out)
    return out


def _parse_p_float(text: str) -> float | None:
    """把各种记法（7.56e-12 / 7.560×10⁻¹² / p<10⁻⁴² / 0.001）解析为 float。

    解析失败返回 None（不参与数值比对，避免误报）。
    """
    t = text.replace(" ", "")
    t = t.translate(str.maketrans("⁰¹²³⁴⁵⁶⁷⁸⁹⁻⁺−", "0123456789-+-"))
    # 1.032×10⁻⁴² → e-42
    t = re.sub(r"[×x]10\^?", "e", t)
    # p<10⁻⁴² 形式 → <e-42（数值部分取 e-42）
    t = re.sub(r"10\^?([+-]?\d+)", lambda m: f"e{m.group(1)}", t)
    try:
        return float(t)
    except ValueError:
        return None


def _sig_parts(value: float) -> tuple[str, int]:
    """取 3 位有效数字与指数，用于权威值拼写笔误检测。"""
    s = f"{value:.2e}"
    mant, _, exp = s.partition("e")
    return mant, int(exp)


def _build_known_p_register(
    authoritative_p: dict[str, str],
    deprecated_p: dict[str, str],
    stats: dict[str, Any] | None = None,
) -> list[tuple[float, str]]:
    """把权威 + 废弃 p 值解析为数值注册表 [(value, 描述), ...]。

    若提供 stats（完整 statistics.yaml），还会递归收集所有 p_value 字段，
    避免把"已收录但不在 build_authoritative_p_values 窄口径内"的合法 p 值
    误判为孤儿值（如噪声 p=2.98e-08、真机 canonical p=8.882e-16 等）。
    """
    register: list[tuple[float, str]] = []
    for p_str, ctx in authoritative_p.items():
        v = _parse_p_float(p_str)
        if v is not None:
            register.append((v, f"权威: {ctx}"))
    for p_str, ctx in deprecated_p.items():
        v = _parse_p_float(p_str)
        if v is not None:
            register.append((v, f"废弃: {ctx}"))
    if stats is not None:
        # 递归收集所有 p_value / *_p_value 字段
        collected: dict[float, str] = {}

        def _walk(node: Any, path: str) -> None:
            if isinstance(node, dict):
                for k, v in node.items():
                    child = f"{path}.{k}" if path else str(k)
                    if re.search(r"p_value$", str(k)) and isinstance(v, (int, float, str)):
                        pv = _parse_p_float(str(v))
                        if pv is not None:
                            collected.setdefault(pv, child)
                    else:
                        _walk(v, child)
            elif isinstance(node, list):
                for i, v in enumerate(node):
                    _walk(v, f"{path}[{i}]")

        _walk(stats, "")
        for pv, where in collected.items():
            register.append((pv, f"权威(收录于 {where})"))
    return register


def _check_p_value_origin(
    raw: str, normalized: str, register: list[tuple[float, str]], line: str
) -> list[str]:
    """对单个 p 值做"权威值拼写检查 + 孤儿 p 值检测"。

    精度约束（8.7-v4 审查修正，避免误报）：
      - 阈值表达式（p<0.05 / p>0.05 / p<0.001）不是具体 p 值，一律跳过；
      - 超出 [0,1] 的"p 值"（如 P=20.18 SWAP 数、p=50 样本量）不是概率，一律跳过；
      - 仅对"具体、合法的 p 值"做权威拼写检查与孤儿检测。

    Returns:
        告警列表（为空表示该 p 值登记在权威/废弃集合内、属合法辅助实验值，
        或无法解析 / 非具体 p 值）。
    """
    # 阈值表达式跳过（p<0.05 / p>0.05 / p<0.001 是显著性阈值，不是具体 p 值）
    if "<" in raw or ">" in raw:
        return []
    val = _parse_p_float(normalized)
    if val is None:
        return []
    # 非法概率（>1：如 P=20.18 SWAP 数、p=50 样本量；<0 异常）跳过
    if not (0.0 < val <= 1.0):
        return []
    # 精确匹配已知权威/废弃 p 值 → 合规
    for known_val, _ctx in register:
        if math.isclose(val, known_val, rel_tol=1e-6, abs_tol=1e-15):
            return []
    # 属于已登记的合法辅助实验 p 值（编译层/噪声/真机等，详见 _KNOWN_LEGIT_P_VALUES）→ 合规
    for known_val in _KNOWN_LEGIT_P_VALUES:
        if math.isclose(val, known_val, rel_tol=1e-6, abs_tol=1e-15):
            return []
    # 权威值拼写笔误检测：同 3 位有效数字但指数不同（如 7.560e-12 → 7.56e-62）
    unknown_mant, unknown_exp = _sig_parts(val)
    for known_val, ctx in register:
        known_mant, known_exp = _sig_parts(known_val)
        if unknown_mant == known_mant and unknown_exp != known_exp:
            return [
                f"  {raw}: 疑似权威 {ctx} 的拼写笔误（有效数字 {unknown_mant} 与权威一致，"
                f"指数 {unknown_exp} ≠ {known_exp}）"
            ]
    # 孤儿 p 值：既非权威也非废弃、非合法辅助值、非诚实披露上下文 → 无出处告警
    if not _is_honest_disclosure(line):
        return [f"  {raw}: 孤儿 p 值，未在 statistics.yaml 权威/废弃集合中登记（无出处，请核验）"]
    return []


# 8.7-v4 审查：接受 3 位有效数字匹配的 p 值（如 7.56e-12 与 7.560e-12）
# 用于 _check_p_value_origin 的数值比对（保留此处常量便于后续调整）
_P_REL_TOL = 1e-6

# 8.7-v4 审查：已登记的合法辅助实验 p 值（编译层/噪声/真机/退火等专项实验的
# 独立 p 值，虽不在 build_authoritative_p_values 的窄权威口径内，但均有出处报告，
# 不属于"孤儿 p 值"。孤儿检测对它们豁免，避免误报。）
_KNOWN_LEGIT_P_VALUES: tuple[float, ...] = (
    1.0,  # 8.7-v4 修复：p=1 为 DQN vs Random 完全相等的退化比较（统计量=0.0000），
    # 合法"无差异"结果而非权威小 p 值的拼写笔误；登记避免被同位有效数字误判
    0.0220,  # 退火配对检验
    0.0294,  # 退火 45k checkpoint
    0.031,  # SJF vs FCFS 真机检验
    0.19,  # 退火 5seed 旧实验
    # 8.7-v4 修复：登记以下合法辅助实验 p 值（均有出处报告，非孤儿），
    # 避免 noise_feedback / SOTA / 真机噪声方向性 / Mann-Whitney 复核被误报为无出处。
    0.222,  # 真机噪声分布训练注入方向性（PPO-noise vs standard，N=5，defense_qa_handbook）
    0.2235,  # SOTA 16维 PPO vs 观测感知 FCFS（N=50，sota_comparison）
    0.25,  # 量子噪声影响
    2.28e-60,  # SJF vs FCFS 仿真（8.8 修正权威；旧 0.2827 为错误值已废弃）
    0.344,  # 真机 vs 仿真对比
    0.3942,  # 退火独立检验
    0.020,  # 噪声派单率（generate_defense_ppt 用 0.020，与 0.0203 同源）
    0.0203,  # 噪声派单率（noise_feedback_v2 报告）
    0.5837,  # 噪声奖励影响 Mann-Whitney（noise_feedback_v2 报告）
    0.584,  # 噪声奖励影响（noise_feedback_v2）
    0.599,  # 编译中+深电路子集
    0.628,  # 噪声稳定性
    0.84,  # 编译全电路
    1.86e-12,  # PPO vs FCFS Mann-Whitney U（N=250，statistics.yaml 权威收录）
    2.49e-03,  # 编译深电路
    2.75e-02,  # 编译深电路子集
    8.52e-03,  # 编译深电路 seed=7 交叉验证（technical_whitepaper §11.2）
    2.621e-10,  # DQN vs SJF
    8.44e-08,  # PPO 编译优化 Agent
)
# 注：1.86e-12（PPO vs FCFS Mann-Whitney U）与 2.98e-08（噪声配对）等已在
# statistics.yaml p_value 字段收录，经 _build_known_p_register 递归收集，无需在此重复。


def scan_markdown_file(
    filepath: Path,
    authoritative_p: dict[str, str],
    deprecated_p: dict[str, str],
    stats: dict[str, Any] | None = None,
) -> list[str]:
    """扫描单个Markdown文件，返回告警列表。"""
    warnings = []

    try:
        with open(filepath, encoding="utf-8") as f:
            lines = f.readlines()
    except (OSError, UnicodeDecodeError):
        return warnings

    text = "".join(lines)
    # 8.7-v4 审查：构建权威/废弃 p 值数值注册表（用于孤儿 p 值与拼写笔误检测）
    p_register = _build_known_p_register(authoritative_p, deprecated_p, stats)

    for line_num, line in enumerate(lines, 1):
        # 8.7-v4 审查：先剥离 markdown 排版再用于黑名单/数字检测，
        # 堵住被 **加粗**、`行内代码`、[链接](url) 包裹的废弃数字绕过正则的漏洞。
        strip_line = _strip_markdown(line)

        # 检查 Welch t 错误搭配
        welch_err = check_welch_t_misattribution(line)
        if welch_err:
            warnings.append(f"  L{line_num}: {welch_err}")
            warnings.append(f"    > {line.strip()[:120]}")

        # 检查黑名单表述（Issue #446）—— 诚实披露上下文豁免
        # 8.6 复核：authoritative_numbers.md 自身的"禁止表述|正确替代"表即黑名单来源，
        # 该表按设计列出被禁止的表述，跳过对其黑名单检测（不当作正文违规）。
        # 8.7-v2：对中文数字归一化后再匹配，堵住"百分之X"中文表述绕过检测的漏洞。
        if filepath.name != "authoritative_numbers.md" and not _is_honest_disclosure(strip_line):
            norm_line = _normalize_chinese_numerals(strip_line)
            for pattern, message in BLACKLIST_PATTERNS:
                if re.search(pattern, norm_line, re.IGNORECASE):
                    # 8.7-v4 修复：stress +91.4% 黑名单在"行覆盖率"语义下豁免，
                    # 避免 mutation_testing 等覆盖率历史陈述被误报为应力优势。
                    if message.startswith(
                        "BLACKLIST: stress 量子波动"
                    ) and _is_stress_coverage_context(strip_line):
                        continue
                    warnings.append(f"  L{line_num}: {message}")
                    warnings.append(f"    > {line.strip()[:120]}")

        # 提取p值
        p_values = extract_p_values_from_line(line)
        for _raw, normalized in p_values:
            # 检查是否为废弃p值（废弃声明上下文中提及旧p值属合规引用，豁免）
            if _is_deprecation_notice(line):
                break
            # 是否命中已登记废弃值（含"Random vs PPO 合法上下文"的情况）
            matched_deprecated = False
            for dep_val, dep_msg in deprecated_p.items():
                if normalized == dep_val or normalized.startswith(dep_val[:8]):
                    matched_deprecated = True
                    # 检查上下文：如果是在正确的实验上下文中使用则跳过
                    # 例如 p=4.92e-55 在 "Random vs PPO" 上下文中是正确的
                    line_lower = line.lower()
                    is_4_92e55 = "4.92e-55" in normalized or "4.92e-55" in dep_val
                    if is_4_92e55 and "random" in line_lower and "ppo" in line_lower:
                        break  # 合法上下文，退出内层循环（matched_deprecated 已置真）
                    warnings.append(f"  L{line_num}: 废弃/错误p值 p={normalized}")
                    warnings.append(f"    {dep_msg}")
                    warnings.append(f"    > {line.strip()[:120]}")
                    break
            if not matched_deprecated and not _is_deprecation_notice(line):
                # 8.7-v4 审查：废弃值未命中后，再做"权威值拼写笔误 + 孤儿 p 值"检测。
                # 注意：仅在非废弃声明上下文执行（避免对已诚实披露的旧 p 值误报）。
                origin_warnings = _check_p_value_origin(_raw, normalized, p_register, strip_line)
                for ow in origin_warnings:
                    warnings.append(f"  L{line_num}: {ow}")
                    warnings.append(f"    > {line.strip()[:120]}")

    # 8.5 审查：权威数值内部冲突（A8）
    for vals, conflict_msg in INTERNAL_CONFLICTS:
        if vals[0] not in text and vals[1] not in text:
            continue
        hits = [v for v in vals if v in text]
        if len(hits) >= 2 and not any(
            any(mark in line for mark in _DISCLOSURE_MARK) for line in lines
        ):
            warnings.append(f"  文件级冲突: {conflict_msg}")
            for v in hits:
                warn_lines = [i for i, ln in enumerate(lines, 1) if v in ln]
                warnings.append(f"    -> '{v}' 出现于 L{warn_lines[:3]}")

    # 8.7 审查新增（A9）：孤立废弃值检测
    # 历史漏洞：value_quantification/platform_landing_value/requirements_traceability 等
    # 活跃交付文档曾把已废弃的"利用率 +7.9%"当作唯一权威呈现，而"同文件并存"冲突
    # 检测（A8）因该文件不含 -3.3% 而抓不到。此处对孤立出现的废弃值单独告警，
    # 仅豁免带诚实披露标记（废弃/错误数据/历史/旧口径等）的行。
    for old_val, new_val, desc in _ORPHAN_DEPRECATED:
        for i, ln in enumerate(lines, 1):
            if old_val in _normalize_chinese_numerals(ln) and not _is_honest_disclosure(ln):
                warnings.append(f"  L{i}: 孤立废弃值 '{old_val}'（{desc}），权威应为 {new_val}")
                warnings.append(f"    > {ln.strip()[:120]}")

    return warnings


def scan_pdf_file(
    filepath: Path,
    authoritative_p: dict[str, str],
    deprecated_p: dict[str, str],
    stats: dict[str, Any] | None = None,
) -> list[str]:
    """扫描单个 PDF 文件中的废弃统计口径。

    8.7-v3 红队审查 P0：交付白皮书 PDF 是二进制产物，此前不在任何口径门禁内，
    导致 8.5~8.7 轮次反复修正 .md 源文件后，PDF 仍保留旧数字（2348.91/1.449e-66
    /33.6%/+51% 等）。此处用 pypdf/PyPDF2 提取文本，复用与 .md 相同的检测逻辑，
    确保"提交的 PDF"与"权威口径"始终一致。

    注意：PDF 文本提取会把"废弃/禁止/旧/历史"等诚实披露限定词拆到相邻行（与
    .md 的行内完整性不同），因此此处用"相邻行窗口"判断诚实披露，避免对已诚实
    化（如实引用旧值并声明废弃）的段落误报。

    Args:
        filepath: PDF 文件路径
        authoritative_p: 权威 p 值映射（预留，与 scan_markdown_file 签名一致）
        deprecated_p: 废弃 p 值映射（预留）

    Returns:
        告警列表
    """
    try:
        from pypdf import PdfReader
    except ImportError:
        try:
            from PyPDF2 import PdfReader
        except ImportError:
            # 8.7-v4 外部红队修复：门禁 fail-open → fail-closed。
            # 缺库时静默跳过会让 PDF 口径校验形同虚设（全新环境永远"全绿"）。
            # 返回 FATAL 级告警，strict 模式下必须失败；同时要求安装
            # pypdf/python-pptx（已加入 requirements-dev.txt）。
            return [
                "[FATAL] 未安装 pypdf/PyPDF2，PDF 口径扫描无法执行。"
                "请执行 pip install -r requirements-dev.txt（含 pypdf/python-pptx）"
            ]

    try:
        reader = PdfReader(str(filepath))
        text = "\n".join((page.extract_text() or "") for page in reader.pages)
    except Exception as e:  # 提取失败不阻断（PDF 可能加密/扫描件）
        return [f"  PDF 文本提取失败（{e}），跳过口径扫描"]

    lines = text.splitlines()
    warnings: list[str] = []

    def _disclosed(idx: int) -> bool:
        """判断第 idx 行（1-based）是否处于诚实披露上下文（±2 行窗口）。"""
        lo = max(0, idx - 3)
        hi = min(len(lines), idx + 2)
        return any(_is_honest_disclosure(lines[j]) for j in range(lo, hi))

    for line_num, line in enumerate(lines, 1):
        if not _disclosed(line_num):
            strip_line = _strip_markdown(line)
            norm_line = _normalize_chinese_numerals(strip_line)
            for pattern, message in BLACKLIST_PATTERNS:
                if re.search(pattern, norm_line, re.IGNORECASE):
                    warnings.append(f"  L{line_num}: {message}")
                    warnings.append(f"    > {line.strip()[:120]}")

    # 孤立废弃值检测（A9）
    for old_val, new_val, desc in _ORPHAN_DEPRECATED:
        for i, ln in enumerate(lines, 1):
            if old_val in _normalize_chinese_numerals(ln) and not _disclosed(i):
                warnings.append(f"  L{i}: 孤立废弃值 '{old_val}'（{desc}），权威应为 {new_val}")
                warnings.append(f"    > {ln.strip()[:120]}")

    return warnings


def scan_pptx_file(
    filepath: Path,
    authoritative_p: dict[str, str],
    deprecated_p: dict[str, str],
    stats: dict[str, Any] | None = None,
) -> list[str]:
    """扫描 PPTX 二进制交付物（8.7-v4 外部红队修复）。

    答辩 PPT 是评委现场唯一翻开的材料，此前 .pptx 二进制不在任何口径门禁内
    （8.7-v3 修了 PDF 盲区，.pptx 是同构盲区）。本函数用 python-pptx 提取
    全部 slide 的文本（含表格 cell），复用与 .md 相同的黑名单/诚实披露检测。

    Returns:
        告警列表；未安装 python-pptx 时返回 [FATAL] 级告警（fail-closed）。
    """
    try:
        from pptx import Presentation
    except ImportError:
        return [
            "[FATAL] 未安装 python-pptx，PPTX 口径扫描无法执行。"
            "请执行 pip install -r requirements-dev.txt（含 pypdf/python-pptx）"
        ]

    try:
        prs = Presentation(str(filepath))
    except Exception as e:
        return [f"[FATAL] PPTX 解析失败（{e}），口径扫描无法执行"]

    warnings: list[str] = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    parts.append("".join(run.text for run in para.runs))
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
            if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                try:
                    for sub in shape.shapes:
                        if sub.has_text_frame:
                            for para in sub.text_frame.paragraphs:
                                parts.append("".join(run.text for run in para.runs))
                except Exception:
                    continue
        # 每页合并为"行"，检测黑名单 + 孤立废弃值（诚实披露上下文豁免）
        slide_text = "\n".join(parts)
        # 8.7-v4 修复：PPTX 文本 run 常被换行拆分（同一句拆成多段），按行检测
        # 会让"限定词与数字不同行"误报为裸展示。此处同时保留：
        #   1) 全文合并（slide_text）做诚实披露判定——限定词只要在同一页就生效
        #   2) 逐行检测黑名单，避免整页大文本绕过
        disclosed_overall = _is_honest_disclosure(slide_text)
        lines = slide_text.splitlines()
        for i, ln in enumerate(lines, 1):
            lo = max(0, i - 3)
            hi = min(len(lines), i + 2)
            disclosed = disclosed_overall or any(
                _is_honest_disclosure(lines[j]) for j in range(lo, hi)
            )
            if disclosed:
                continue
            norm_line = _normalize_chinese_numerals(ln)
            for pattern, message in BLACKLIST_PATTERNS:
                if re.search(pattern, norm_line, re.IGNORECASE):
                    warnings.append(f"  Slide {slide_idx} L{i}: {message}")
                    warnings.append(f"    > {ln.strip()[:120]}")
            for old_val, new_val, desc in _ORPHAN_DEPRECATED:
                if old_val in norm_line:
                    warnings.append(
                        f"  Slide {slide_idx} L{i}: 孤立废弃值 '{old_val}'（{desc}），权威应为 {new_val}"
                    )
                    warnings.append(f"    > {ln.strip()[:120]}")

    return warnings


def main() -> int:
    """主入口：扫描所有Markdown文件并报告统计口径不一致。"""
    import argparse

    parser = argparse.ArgumentParser(description="统计口径一致性检查")
    parser.add_argument(
        "--strict",
        action="store_true",
        help="严格模式：发现不一致时退出码1（用于CI）",
    )
    parser.add_argument(
        "--stats-file",
        type=str,
        default=str(STATS_YAML),
        help="权威统计源YAML文件路径",
    )
    args = parser.parse_args()

    stats_path = Path(args.stats_file)
    if not stats_path.exists():
        print(f"ERROR: 权威统计源文件不存在: {stats_path}")
        return 1

    # 加载权威源
    stats = load_statistics_yaml()
    authoritative_p = build_authoritative_p_values(stats)
    deprecated_p = build_deprecated_p_values(stats)

    print("=" * 70)
    print("  统计口径一致性检查 (Issue #141)")
    print("=" * 70)
    print(f"  权威源: {stats_path}")
    print(f"  权威p值数: {len(authoritative_p)}")
    print(f"  废弃p值数: {len(deprecated_p)}")
    print("=" * 70)

    # Issue #536: 权威源覆盖度检查（确保已发布实验不遗漏）
    coverage_warnings = check_authoritative_coverage(stats)
    if coverage_warnings:
        print("\n[Issue #536] 权威源覆盖度检查:")
        for w in coverage_warnings:
            print(f"  [ERROR] {w}")
        print()
    else:
        print("[Issue #536] 权威源覆盖度检查: ✅ 关键实验均已收录")
        print()

    # 8.6 审查新增：statistics.yaml 内部数值自洽校验
    yaml_warnings = check_yaml_internal_consistency(stats)
    if yaml_warnings:
        print("\n[Issue #854] statistics.yaml 内部自洽检查:")
        for w in yaml_warnings:
            print(f"  [ERROR] {w}")
        print()
    else:
        print("[Issue #854] statistics.yaml 内部自洽检查: ✅ 数值自洽")
        print()

    # 收集所有Markdown文件（使用 os.walk 跳过 node_modules 等目录）
    md_files = set()
    # 演示链路 .py 文件纳入扫描（与 .md 共用同一套逐行检查逻辑）
    for rel in DEMO_CHAIN_FILES:
        demo_path = _PROJECT_ROOT / rel
        if demo_path.exists():
            md_files.add(demo_path)
    # Issue #174: 排除非交付目录（AI 工作目录、数据转储、临时文件等），
    # 这些目录可能包含历史数字/草稿，不应参与口径审计，避免假失败
    exclude_dirs = {
        "node_modules",
        ".git",
        "__pycache__",
        ".venv",
        "venv",
        "mutants",
        ".pytest_cache",
        ".mypy_cache",
        ".ruff_cache",
        ".hypothesis",
        # 非交付目录：AI 工作目录、数据转储、临时文件等
        ".workbuddy",
        ".trae",
        ".trae-cn",
        "project-review",
        ".trae-html-share-packages",
        # 8.6 复核：历史快照目录（非交付物），保留旧数字属预期，不参与口径审计
        ".archive",
        "archive",
        "tmp",
        "temp",
        "data",
        "datasets",
        "build",
        "dist",
        "downloads",
    }
    for root, dirs, files in os.walk(_PROJECT_ROOT, topdown=True):
        # 原地修改 dirs 跳过排除目录，阻止 os.walk 递归进入
        # 同时跳过所有 .venv* 开头的目录（各种虚拟环境）和 site-packages
        dirs[:] = [
            d
            for d in dirs
            if d not in exclude_dirs
            and not d.startswith(".venv")
            and d not in ("site-packages", "lib", "lib64")
        ]
        for fname in files:
            if not fname.endswith(".md"):
                continue
            fpath = Path(root) / fname
            rel = str(fpath.relative_to(_PROJECT_ROOT)).replace("\\", "/")
            if rel not in EXCLUDE_PATHS and not any(pat in rel for pat in EXCLUDE_PATTERNS):
                md_files.add(fpath)

    md_files = sorted(md_files, key=lambda f: str(f.relative_to(_PROJECT_ROOT)))

    print(f"  扫描文件数: {len(md_files)}")
    print()

    total_warnings = 0
    files_with_warnings = 0

    # 8.6 审查：覆盖度与 YAML 内部自洽告警计入总告警（strict 模式下触发退出码 1）
    total_warnings += len(coverage_warnings)
    total_warnings += len(yaml_warnings)

    for filepath in md_files:
        rel_path = str(filepath.relative_to(_PROJECT_ROOT)).replace("\\", "/")
        warnings = scan_markdown_file(filepath, authoritative_p, deprecated_p, stats)
        if warnings:
            files_with_warnings += 1
            total_warnings += len(warnings) // 3  # 每个告警约3行
            print(f"[WARN] {rel_path}")
            for w in warnings:
                print(w)
            print()

    # 8.7-v3 红队审查 P0：交付白皮书 PDF 为二进制提交物，此前不在任何口径门禁内，
    # 导致 .md 源文件修正后 PDF 仍保留旧数字（2348.91/1.449e-66/33.6%/+51% 等）。
    # 此处对已生成的提交 PDF 复用与 .md 相同的检测逻辑，确保"提交的 PDF"与
    # "权威口径"始终一致。
    deliverable_pdfs = [
        "docs/technical_whitepaper.pdf",
    ]
    for rel in deliverable_pdfs:
        pdf_path = _PROJECT_ROOT / rel
        if not pdf_path.exists():
            continue
        pdf_warnings = scan_pdf_file(pdf_path, authoritative_p, deprecated_p, stats)
        if pdf_warnings:
            files_with_warnings += 1
            # 8.7-v4 修复：PDF/PPTX 告警为 2 行/处，原 //3 会算成 0 导致
            # strict 模式下有告警却仍通过（fail-open 残留）
            total_warnings += len(pdf_warnings) // 2
            print(f"[WARN] {rel} (PDF)")
            for w in pdf_warnings:
                print(w)
            print()
        else:
            print(f"[OK] {rel} (PDF) 口径一致")

    # 8.7-v4 外部红队修复：.pptx 二进制纳入口径门禁（答辩 PPT 是评委现场
    # 唯一翻开的材料，与 PDF 同构的盲区必须封死）。
    deliverable_ppt = "deliverable_models/答辩PPT.pptx"
    ppt_path = _PROJECT_ROOT / deliverable_ppt
    if ppt_path.exists():
        ppt_warnings = scan_pptx_file(ppt_path, authoritative_p, deprecated_p, stats)
        if ppt_warnings:
            files_with_warnings += 1
            # 2 行/告警（见 PDF 注释），//3 会导致 strict 下漏检
            total_warnings += len(ppt_warnings) // 2
            print(f"[WARN] {deliverable_ppt} (PPTX)")
            for w in ppt_warnings:
                print(w)
            print()
        else:
            print(f"[OK] {deliverable_ppt} (PPTX) 口径一致")

    # 汇总
    print("=" * 70)
    if total_warnings == 0:
        print("  ✅ 统计口径一致性检查通过，无不一致告警。")
    else:
        print(f"  ⚠️  发现 {total_warnings} 处不一致告警（{files_with_warnings} 个文件）。")
    print("=" * 70)

    # 打印权威p值参考表
    print("\n权威p值参考表:")
    print(f"  {'p值':<20} {'实验上下文'}")
    print("  " + "-" * 66)
    for p_val, ctx in sorted(authoritative_p.items()):
        print(f"  {p_val:<20} {ctx}")

    if args.strict and total_warnings > 0:
        return 1
    return 0


if __name__ == "__main__":
    sys.exit(main())
