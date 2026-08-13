#!/usr/bin/env python3
"""
M5 最终提交物一键打包与版本校验脚本

功能：
1. --check 模式：校验所有提交物是否符合清单要求
2. --pack 模式：校验 + 创建最终提交压缩包
3. --report PATH 模式：将校验结果输出为 Markdown 格式的缺失项清单
4. --prepare 模式：准备可自动生成的提交物（dist 目录、缺失项报告、检查清单）

作者：量子RL调度系统团队
日期：2026-07-02
"""

import argparse
import json
import shutil
import subprocess
import sys
import zipfile
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any, ClassVar

import yaml

# 修复 Windows GBK 终端下 emoji 字符导致的 UnicodeEncodeError 崩溃
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")


@dataclass
class ItemResult:
    """单个提交物的校验结果"""

    item_id: str
    name: str
    item_type: str
    path: str
    passed: bool
    messages: list[str] = field(default_factory=list)
    severity: str = "error"  # "error" 或 "warning"


class SubmissionValidator:
    """提交物校验器"""

    def __init__(
        self,
        manifest_path: str,
        project_root: str = ".",
        skip_items: list[str] | None = None,
    ) -> None:
        """初始化校验器

        Args:
            manifest_path: 清单文件路径
            project_root: 项目根目录
            skip_items: 需要跳过的提交物 id 列表（如 release CI 中跳过仅
                本地生成的交付物：dist zip / 演示视频等，Issue #860）
        """
        with open(manifest_path, encoding="utf-8") as f:
            self.manifest = yaml.safe_load(f)
        self.project_root = Path(project_root)
        self.skip_items: set[str] = set(skip_items or [])
        self.errors: list[str] = []
        self.warnings: list[str] = []
        self.results: list[ItemResult] = []

    def validate_all(self) -> bool:
        """校验所有提交物

        Returns:
            是否通过校验
        """
        print("=== M5 提交物校验报告 ===")
        print(f"时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        print(f"版本: {self.manifest['submission']['version']}")
        print(f"截止日期: {self.manifest['submission']['deadline']}")
        print()

        for item in self.manifest["items"]:
            if item["id"] in self.skip_items:
                print(f"[SKIP] {item['id']}: 已按 --skip-items 跳过（本地人工交付物）")
                continue
            self._validate_item(item)

        self._check_version_consistency()
        self._report()
        return len(self.errors) == 0

    def _validate_item(self, item: dict[str, Any]) -> None:
        """校验单个提交物

        Args:
            item: 提交物定义
        """
        item_id = item["id"]
        item_name = item["name"]
        item_type = item["type"]
        path = self.project_root / item["path"]
        messages: list[str] = []

        print(f"[{item_id}] {item_name} ({item_type})")

        # 检查文件存在性
        if not path.exists():
            # 对白皮书特殊处理：manifest 要求 pdf，但可能存在 docx 源文件
            if item_type == "pdf":
                docx_path = path.with_suffix(".docx")
                if docx_path.exists():
                    msg = (
                        f"文件不存在: {path}，但发现 docx 源文件: {docx_path.name}，"
                        f"需转换为 PDF 后再提交"
                    )
                    self.warnings.append(f"[{item_id}] {msg}")
                    print(f"  ⚠️  {msg}")
                    messages.append(msg)
                    self.results.append(
                        ItemResult(
                            item_id=item_id,
                            name=item_name,
                            item_type=item_type,
                            path=str(item["path"]),
                            passed=False,
                            messages=messages,
                            severity="warning",
                        )
                    )
                    return
            self.errors.append(f"[{item_id}] 文件不存在: {path}")
            print(f"  ❌ 文件不存在: {path}")
            messages.append(f"文件不存在: {path}")
            self.results.append(
                ItemResult(
                    item_id=item_id,
                    name=item_name,
                    item_type=item_type,
                    path=str(item["path"]),
                    passed=False,
                    messages=messages,
                    severity="error",
                )
            )
            return

        # 按类型校验
        errors_before = len(self.errors)
        if item_type == "pdf":
            self._validate_pdf(item, path, messages)
        elif item_type == "pptx":
            self._validate_pptx(item, path, messages)
        elif item_type == "mp4":
            self._validate_mp4(item, path, messages)
        elif item_type == "zip":
            self._validate_zip(item, path, messages)
        elif item_type == "git_tag":
            self._validate_git_tag(item, messages)
        elif item_type == "md":
            self._validate_markdown(item, path, messages)
        elif item_type == "directory":
            self._validate_directory(item, path, messages)

        # 检查依赖
        if "depends_on" in item:
            self._check_dependency(item)

        # 记录结果（仅当未被前面的提前 return 记录过时）
        has_error = len(self.errors) > errors_before
        self.results.append(
            ItemResult(
                item_id=item_id,
                name=item_name,
                item_type=item_type,
                path=str(item["path"]),
                passed=not has_error,
                messages=messages,
                severity="error" if has_error else "info",
            )
        )

    def _validate_pdf(self, item: dict[str, Any], path: Path, messages: list[str]) -> None:
        """校验 PDF 文件

        Args:
            item: 提交物定义
            path: 文件路径
            messages: 用于收集本项校验消息的列表
        """
        try:
            from PyPDF2 import PdfReader

            reader = PdfReader(path)
            num_pages = len(reader.pages)
            reqs = item.get("requirements", {})

            min_pages = reqs.get("min_pages")
            max_pages = reqs.get("max_pages")

            if min_pages and num_pages < min_pages:
                msg = f"PDF 页数不足: {num_pages} < {min_pages}"
                self.errors.append(f"[{item['id']}] {msg}")
                messages.append(msg)
                print(f"  ❌ {msg}")
            elif max_pages and num_pages > max_pages:
                msg = f"PDF 页数超限: {num_pages} > {max_pages}"
                self.errors.append(f"[{item['id']}] {msg}")
                messages.append(msg)
                print(f"  ❌ {msg}")
            else:
                messages.append(f"页数: {num_pages}")
                print(f"  ✅ 页数: {num_pages}")

            # 检查必需内容
            must_contain = reqs.get("must_contain", [])
            if must_contain:
                text = ""
                for page in reader.pages:
                    text += page.extract_text()

                missing = [kw for kw in must_contain if kw not in text]
                if missing:
                    msg = f"PDF 缺少关键词: {', '.join(missing)}"
                    self.warnings.append(f"[{item['id']}] {msg}")
                    messages.append(msg)
                    print(f"  ⚠️  {msg}")
                else:
                    messages.append("包含所有必需关键词")
                    print("  ✅ 包含所有必需关键词")

        except ImportError:
            msg = "PyPDF2 未安装，跳过 PDF 详细校验"
            self.warnings.append(f"[{item['id']}] {msg}")
            messages.append(msg)
            print(f"  ⚠️  {msg}")
        except Exception as e:
            msg = f"PDF 校验失败: {e}"
            self.errors.append(f"[{item['id']}] {msg}")
            messages.append(msg)
            print(f"  ❌ {msg}")

    def _validate_pptx(self, item: dict[str, Any], path: Path, messages: list[str]) -> None:
        """校验 PPTX 文件

        Args:
            item: 提交物定义
            path: 文件路径
            messages: 用于收集本项校验消息的列表
        """
        try:
            from pptx import Presentation

            prs = Presentation(path)
            num_slides = len(prs.slides)
            reqs = item.get("requirements", {})

            min_slides = reqs.get("min_slides")
            max_slides = reqs.get("max_slides")

            if min_slides and num_slides < min_slides:
                msg = f"PPT 页数不足: {num_slides} < {min_slides}"
                self.errors.append(f"[{item['id']}] {msg}")
                messages.append(msg)
                print(f"  ❌ {msg}")
            elif max_slides and num_slides > max_slides:
                msg = f"PPT 页数超限: {num_slides} > {max_slides}"
                self.errors.append(f"[{item['id']}] {msg}")
                messages.append(msg)
                print(f"  ❌ {msg}")
            else:
                messages.append(f"幻灯片数: {num_slides}")
                print(f"  ✅ 幻灯片数: {num_slides}")

            # 检查必需幻灯片
            must_contain = reqs.get("must_contain_slides", [])
            if must_contain:
                slide_titles = []
                for slide in prs.slides:
                    if slide.shapes.title:
                        slide_titles.append(slide.shapes.title.text)
                    # 无正式 title 占位符的幻灯片：收集全部文本作为标题候选，
                    # 避免"标题在文本框中"的常见版式被误判为缺少必需页
                    else:
                        texts = []
                        for shape in slide.shapes:
                            if shape.has_text_frame and shape.text_frame.text.strip():
                                texts.append(shape.text_frame.text)
                        if texts:
                            slide_titles.append("\n".join(texts))

                missing = [
                    title for title in must_contain if not any(title in t for t in slide_titles)
                ]
                if missing:
                    msg = f"PPT 缺少幻灯片: {', '.join(missing)}"
                    self.warnings.append(f"[{item['id']}] {msg}")
                    messages.append(msg)
                    print(f"  ⚠️  {msg}")
                else:
                    messages.append("包含所有必需幻灯片")
                    print("  ✅ 包含所有必需幻灯片")

        except ImportError:
            msg = "python-pptx 未安装，跳过 PPTX 详细校验"
            self.warnings.append(f"[{item['id']}] {msg}")
            messages.append(msg)
            print(f"  ⚠️  {msg}")
        except Exception as e:
            msg = f"PPTX 校验失败: {e}"
            self.errors.append(f"[{item['id']}] {msg}")
            messages.append(msg)
            print(f"  ❌ {msg}")

    def _validate_mp4(self, item: dict[str, Any], path: Path, messages: list[str]) -> None:
        """校验 MP4 文件

        Args:
            item: 提交物定义
            path: 文件路径
            messages: 用于收集本项校验消息的列表
        """
        try:
            # 检查文件大小
            size_mb = path.stat().st_size / (1024 * 1024)
            reqs = item.get("requirements", {})
            max_size = reqs.get("max_size_mb")

            if max_size and size_mb > max_size:
                msg = f"视频文件过大: {size_mb:.1f}MB > {max_size}MB"
                self.errors.append(f"[{item['id']}] {msg}")
                messages.append(msg)
                print(f"  ❌ {msg}")
            else:
                messages.append(f"文件大小: {size_mb:.1f}MB")
                print(f"  ✅ 文件大小: {size_mb:.1f}MB")

            # 使用 ffprobe 检查时长和分辨率
            result = subprocess.run(
                [
                    "ffprobe",
                    "-v",
                    "error",
                    "-select_streams",
                    "v:0",
                    "-show_entries",
                    "stream=width,height,duration",
                    "-show_entries",
                    "format=duration",
                    "-of",
                    "json",
                    str(path),
                ],
                capture_output=True,
                text=True,
                check=True,
            )

            info = json.loads(result.stdout)
            duration = float(info["format"]["duration"])
            stream = info["streams"][0]
            width = stream["width"]
            height = stream["height"]

            min_duration = reqs.get("min_duration_seconds")
            max_duration = reqs.get("max_duration_seconds")

            if min_duration and duration < min_duration:
                msg = f"视频时长不足: {duration:.1f}s < {min_duration}s"
                self.errors.append(f"[{item['id']}] {msg}")
                messages.append(msg)
                print(f"  ❌ {msg}")
            elif max_duration and duration > max_duration:
                msg = f"视频时长超限: {duration:.1f}s > {max_duration}s"
                self.errors.append(f"[{item['id']}] {msg}")
                messages.append(msg)
                print(f"  ❌ {msg}")
            else:
                messages.append(f"时长: {duration:.1f}s")
                print(f"  ✅ 时长: {duration:.1f}s")

            expected_resolution = reqs.get("resolution")
            if expected_resolution:
                exp_w, exp_h = map(int, expected_resolution.split("x"))
                if width != exp_w or height != exp_h:
                    msg = f"视频分辨率不匹配: {width}x{height} != {expected_resolution}"
                    self.errors.append(f"[{item['id']}] {msg}")
                    messages.append(msg)
                    print(f"  ❌ {msg}")
                else:
                    messages.append(f"分辨率: {width}x{height}")
                    print(f"  ✅ 分辨率: {width}x{height}")

        except FileNotFoundError:
            msg = "ffprobe 未安装，跳过 MP4 详细校验"
            self.warnings.append(f"[{item['id']}] {msg}")
            messages.append(msg)
            print(f"  ⚠️  {msg}")
        except subprocess.CalledProcessError as e:
            msg = f"ffprobe 执行失败: {e}"
            self.errors.append(f"[{item['id']}] {msg}")
            messages.append(msg)
            print(f"  ❌ {msg}")
        except Exception as e:
            msg = f"MP4 校验失败: {e}"
            self.errors.append(f"[{item['id']}] {msg}")
            messages.append(msg)
            print(f"  ❌ {msg}")

    def _validate_zip(self, item: dict[str, Any], path: Path, messages: list[str]) -> None:
        """校验 ZIP 文件

        校验维度：
        1. 文件大小不超过 max_size_mb
        2. 必须包含 include 列表中的路径
        3. 不得包含 exclude 列表中的路径

        Args:
            item: 提交物定义
            path: 文件路径
            messages: 用于收集本项校验消息的列表
        """
        size_mb = path.stat().st_size / (1024 * 1024)
        reqs = item.get("requirements", {})
        max_size = reqs.get("max_size_mb")

        if max_size and size_mb > max_size:
            msg = f"ZIP 文件过大: {size_mb:.1f}MB > {max_size}MB"
            self.errors.append(f"[{item['id']}] {msg}")
            messages.append(msg)
            print(f"  ❌ {msg}")
        else:
            messages.append(f"文件大小: {size_mb:.1f}MB")
            print(f"  ✅ 文件大小: {size_mb:.1f}MB")

        # 校验 ZIP 内容（include / exclude 规则）
        include_list = reqs.get("include", [])
        exclude_list = reqs.get("exclude", [])
        if include_list or exclude_list:
            try:
                with zipfile.ZipFile(path, "r") as zipf:
                    zip_names = set(zipf.namelist())
            except zipfile.BadZipFile as e:
                msg = f"ZIP 文件损坏: {e}"
                self.errors.append(f"[{item['id']}] {msg}")
                messages.append(msg)
                print(f"  ❌ {msg}")
                return

            # 检查 include：ZIP 中必须包含这些路径前缀
            if include_list:
                missing_includes = []
                for inc in include_list:
                    # inc 可能是目录（如 "src/"）或文件（如 "README.md"）
                    inc_clean = inc.rstrip("/")
                    found = any(n.startswith(inc_clean + "/") or n == inc_clean for n in zip_names)
                    if not found:
                        missing_includes.append(inc)
                if missing_includes:
                    msg = f"ZIP 缺少必需路径: {', '.join(missing_includes)}"
                    self.errors.append(f"[{item['id']}] {msg}")
                    messages.append(msg)
                    print(f"  ❌ {msg}")
                else:
                    messages.append(f"包含所有必需路径: {', '.join(include_list)}")
                    print("  ✅ 包含所有必需路径")

            # 检查 exclude：ZIP 中不得包含这些路径前缀
            if exclude_list:
                found_excludes = []
                for exc in exclude_list:
                    exc_clean = exc.rstrip("/")
                    found = any(n.startswith(exc_clean + "/") or n == exc_clean for n in zip_names)
                    if found:
                        found_excludes.append(exc)
                if found_excludes:
                    msg = f"ZIP 包含禁止路径: {', '.join(found_excludes)}"
                    self.errors.append(f"[{item['id']}] {msg}")
                    messages.append(msg)
                    print(f"  ❌ {msg}")
                else:
                    messages.append(f"未包含禁止路径: {', '.join(exclude_list)}")
                    print("  ✅ 未包含禁止路径")

    def _validate_git_tag(self, item: dict[str, Any], messages: list[str]) -> None:
        """校验 Git 标签

        除标签存在性外，还会按 requirements 声明校验：
        - ``must_pass_ci``：最新 CI 运行是否通过（Issue #912）
        - ``must_have_readme``：指定目录下是否存在 README.md（Issue #912）

        Args:
            item: 提交物定义
            messages: 用于收集本项校验消息的列表
        """
        reqs = item.get("requirements", {})
        tag = reqs.get("tag")

        if tag:
            try:
                result = subprocess.run(
                    ["git", "tag", "-l", tag],
                    capture_output=True,
                    text=True,
                    check=True,
                    cwd=self.project_root,
                )
                if tag in result.stdout:
                    messages.append(f"标签存在: {tag}")
                    print(f"  ✅ 标签存在: {tag}")
                else:
                    msg = f"Git 标签不存在: {tag}"
                    self.errors.append(f"[{item['id']}] {msg}")
                    messages.append(msg)
                    print(f"  ❌ {msg}")
            except Exception as e:
                msg = f"Git 标签校验失败: {e}"
                self.errors.append(f"[{item['id']}] {msg}")
                messages.append(msg)
                print(f"  ❌ {msg}")

        # 校验最新 CI 运行是否通过（Issue #912）
        if reqs.get("must_pass_ci", False):
            self._check_must_pass_ci(item["id"], messages)

        # 校验 README.md 是否存在（Issue #912）
        if reqs.get("must_have_readme", False):
            self._check_must_have_readme(item["id"], item.get("path", "."), messages)

    def _check_must_pass_ci(self, item_id: str, messages: list[str]) -> None:
        """检查最新一次 CI 运行是否通过

        通过 ``gh run list`` 查询最近一次 GitHub Actions 运行的状态。当 gh CLI
        不可用时仅给出警告（不阻断校验）；当最近一次运行已结束但结论非 success
        时记为错误；运行尚未完成时记为警告。

        Args:
            item_id: 提交物 id
            messages: 用于收集本项校验消息的列表
        """
        gh_path = shutil.which("gh")
        if gh_path is None:
            msg = "gh CLI 不可用，跳过 CI 状态校验"
            self.warnings.append(f"[{item_id}] {msg}")
            messages.append(msg)
            print(f"  ⚠️  {msg}")
            return

        try:
            result = subprocess.run(
                ["gh", "run", "list", "--limit", "1", "--json", "status,conclusion"],
                capture_output=True,
                text=True,
                check=True,
                cwd=self.project_root,
            )
            runs = json.loads(result.stdout)
            if not runs:
                msg = "未找到任何 CI 运行记录，跳过 CI 状态校验"
                self.warnings.append(f"[{item_id}] {msg}")
                messages.append(msg)
                print(f"  ⚠️  {msg}")
                return

            latest = runs[0]
            status = latest.get("status", "")
            conclusion = latest.get("conclusion", "")

            if status != "completed":
                msg = f"最新 CI 运行尚未完成 (status={status})"
                self.warnings.append(f"[{item_id}] {msg}")
                messages.append(msg)
                print(f"  ⚠️  {msg}")
            elif conclusion == "success":
                messages.append(f"最新 CI 运行通过 (conclusion={conclusion})")
                print(f"  ✅ 最新 CI 运行通过 (conclusion={conclusion})")
            else:
                msg = f"最新 CI 运行未通过 (status={status}, conclusion={conclusion})"
                self.errors.append(f"[{item_id}] {msg}")
                messages.append(msg)
                print(f"  ❌ {msg}")
        except subprocess.CalledProcessError as e:
            msg = f"gh run list 执行失败: {e}"
            self.warnings.append(f"[{item_id}] {msg}")
            messages.append(msg)
            print(f"  ⚠️  {msg}")
        except (json.JSONDecodeError, KeyError, TypeError) as e:
            msg = f"gh run list 输出解析失败: {e}"
            self.warnings.append(f"[{item_id}] {msg}")
            messages.append(msg)
            print(f"  ⚠️  {msg}")

    def _check_must_have_readme(self, item_id: str, item_path: str, messages: list[str]) -> None:
        """检查 README.md 是否存在于指定目录

        Args:
            item_id: 提交物 id
            item_path: 提交物所在目录（相对项目根目录，如 "."）
            messages: 用于收集本项校验消息的列表
        """
        readme_path = self.project_root / item_path / "README.md"
        if readme_path.exists():
            messages.append(f"README.md 存在: {readme_path}")
            print(f"  ✅ README.md 存在: {readme_path}")
        else:
            msg = f"README.md 不存在: {readme_path}"
            self.errors.append(f"[{item_id}] {msg}")
            messages.append(msg)
            print(f"  ❌ {msg}")

    def _validate_markdown(self, item: dict[str, Any], path: Path, messages: list[str]) -> None:
        """校验 Markdown 文件

        Args:
            item: 提交物定义
            path: 文件路径
            messages: 用于收集本项校验消息的列表
        """
        if item.get("must_exist", False):
            messages.append("文件存在")
            print("  ✅ 文件存在")

    def _validate_directory(self, item: dict[str, Any], path: Path, messages: list[str]) -> None:
        """校验目录型提交物

        校验维度（8.7-v3 红队审查 P1-4）：
        1. 目录非空（至少含一个文件）
        2. requirements.min_files：至少包含的文件数
        3. requirements.must_contain：必须包含的文件名子串
        4. requirements.exclude_ext：不得包含的扩展名

        Args:
            item: 提交物定义
            path: 文件路径
            messages: 用于收集本项校验消息的列表
        """
        reqs = item.get("requirements", {})

        if not path.is_dir():
            msg = f"不是目录: {path}"
            self.errors.append(f"[{item['id']}] {msg}")
            messages.append(msg)
            print(f"  ❌ {msg}")
            return

        files = [p for p in path.rglob("*") if p.is_file()]
        if not files:
            msg = f"目录为空: {path}"
            self.errors.append(f"[{item['id']}] {msg}")
            messages.append(msg)
            print(f"  ❌ {msg}")
            return

        # 最小文件数校验
        min_files = reqs.get("min_files")
        if min_files and len(files) < min_files:
            msg = f"文件数不足: {len(files)} < {min_files}"
            self.errors.append(f"[{item['id']}] {msg}")
            messages.append(msg)
            print(f"  ❌ {msg}")
        else:
            messages.append(f"文件数: {len(files)}")
            print(f"  ✅ 文件数: {len(files)}")

        # 必需文件名校验（子串匹配）
        must_contain = reqs.get("must_contain", [])
        if must_contain:
            names = [f.name for f in files]
            missing = [kw for kw in must_contain if not any(kw in n for n in names)]
            if missing:
                msg = f"目录缺少必需文件: {', '.join(missing)}"
                self.errors.append(f"[{item['id']}] {msg}")
                messages.append(msg)
                print(f"  ❌ {msg}")
            else:
                messages.append(f"包含必需文件: {', '.join(must_contain)}")
                print(f"  ✅ 包含必需文件: {', '.join(must_contain)}")

        # 禁止扩展名校验
        exclude_ext = reqs.get("exclude_ext", [])
        if exclude_ext:
            forbidden = [f.name for f in files if f.suffix.lower() in exclude_ext]
            if forbidden:
                msg = f"目录包含禁止扩展名文件: {', '.join(forbidden[:5])}"
                self.errors.append(f"[{item['id']}] {msg}")
                messages.append(msg)
                print(f"  ❌ {msg}")
            else:
                messages.append("未包含禁止扩展名文件")
                print("  ✅ 未包含禁止扩展名文件")

    def _check_dependency(self, item: dict[str, Any]) -> None:
        """检查依赖项

        Args:
            item: 提交物定义
        """
        depends_on = item.get("depends_on")
        if depends_on:
            # 简化处理：假设依赖项已满足
            self.warnings.append(f"[{item['id']}] 依赖项: {depends_on}")
            print(f"  ⚠️  依赖项: {depends_on}")

    def _check_version_consistency(self) -> None:
        """检查版本一致性

        校验 manifest 中声明的版本号是否与项目关键文件（README.md、AGENTS.md）
        中出现的版本号一致，避免提交物版本与文档版本不匹配。
        """
        version = self.manifest["submission"]["version"]
        print(f"\n[版本一致性] 目标版本: {version}")

        # 检查 README.md 和 AGENTS.md 中的版本号
        files_to_check = [
            self.project_root / "README.md",
            self.project_root / "AGENTS.md",
        ]

        for file_path in files_to_check:
            if not file_path.exists():
                continue
            try:
                content = file_path.read_text(encoding="utf-8")
                # 检查文件中是否出现版本号
                if version in content:
                    print(f"  ✅ {file_path.name}: 包含版本号 {version}")
                else:
                    msg = f"{file_path.name}: 未找到版本号 {version}"
                    self.warnings.append(f"[版本一致性] {msg}")
                    print(f"  ⚠️  {msg}")
            except OSError as e:
                msg = f"{file_path.name}: 读取失败 ({e})"
                self.warnings.append(f"[版本一致性] {msg}")
                print(f"  ⚠️  {msg}")

    def _report(self) -> None:
        """输出校验报告"""
        print()
        print("=" * 60)
        print(f"错误: {len(self.errors)}")
        for e in self.errors:
            print(f"  [ERROR] {e}")
        print(f"警告: {len(self.warnings)}")
        for w in self.warnings:
            print(f"  [WARN] {w}")
        print("=" * 60)

        if self.errors:
            print("\n❌ 校验失败，存在错误需要修复")
        else:
            print("\n✅ 校验通过，所有提交物符合要求")

    # 缺失项与建议处理方式的映射，用于生成跟踪报告
    MISSING_ITEM_GUIDANCE: ClassVar[dict[str, str]] = {
        "CODE_REPO": "在代码冻结日（2026-08-15）后由管理员执行 `git tag v9.1-submission` 并推送标签",
        "CODE_ARCHIVE": "代码冻结后执行 `python scripts/ci/validate_submission.py --pack` 生成压缩包",
        "WHITEPAPER": "将 `技术白皮书_量子RL调度系统_v3.docx` 导出为 PDF（20-50 页，需含摘要/目录/参考文献）",
        "PRESENTATION": "根据 `答辩PPT大纲.md` 制作 .pptx 文件（15-20 页，需含封面/问题定义/架构图/实验结果/团队介绍）",
        "DEMO_VIDEO": "录制 4-5 分钟 1080p 演示视频（关联 Issue #169）",
        "SUMMARY_REPORT": "按比赛方案第八条第八点生成 `参赛总结报告.pdf`（设计说明/技术实现/结果/创新点总结，PDF 格式）",
        "REGISTRATION_FORM": "提交报名系统中审核通过的 `参赛报名表.pdf`（比赛方案第八条第八点硬性要求）",
    }

    def generate_report(self, output_path: str) -> None:
        """生成 Markdown 格式的缺失项清单报告

        Args:
            output_path: 报告输出路径
        """
        passed_items = [r for r in self.results if r.passed]
        failed_items = [r for r in self.results if not r.passed]
        warning_items = [r for r in self.results if r.severity == "warning"]

        version = self.manifest["submission"]["version"]
        deadline = self.manifest["submission"]["deadline"]

        lines: list[str] = []
        lines.append("# 提交物校验报告 — Issue #168")
        lines.append("")
        lines.append(f"- **版本**: {version}")
        lines.append(f"- **截止日期**: {deadline}")
        lines.append(f"- **生成时间**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        lines.append(
            f"- **总数**: {len(self.results)} 项  |  ✅ 通过: {len(passed_items)}  |  ❌ 缺失: {len(failed_items)}"
        )
        lines.append("")
        lines.append("---")
        lines.append("")

        # 缺失项清单
        lines.append("## ❌ 缺失项清单（需处理）")
        lines.append("")
        if not failed_items:
            lines.append("无缺失项，所有交付物已就位。")
            lines.append("")
        else:
            lines.append("| 编号 | 名称 | 类型 | 期望路径 | 严重度 | 说明 | 建议处理方式 |")
            lines.append("|:--:|:--|:--:|:--|:--:|:--|:--|")
            for r in failed_items:
                guidance = self.MISSING_ITEM_GUIDANCE.get(r.item_id, "—")
                msg_text = "; ".join(r.messages) if r.messages else "—"
                lines.append(
                    f"| {r.item_id} | {r.name} | {r.item_type} | `{r.path}` | {r.severity} | {msg_text} | {guidance} |"
                )
            lines.append("")

        # 警告项清单
        if warning_items:
            lines.append("## ⚠️ 警告项清单（建议关注）")
            lines.append("")
            lines.append("| 编号 | 名称 | 说明 |")
            lines.append("|:--:|:--|:--|")
            for r in warning_items:
                msg_text = "; ".join(r.messages) if r.messages else "—"
                lines.append(f"| {r.item_id} | {r.name} | {msg_text} |")
            lines.append("")

        # 已通过项清单
        lines.append("## ✅ 已通过项清单")
        lines.append("")
        if not passed_items:
            lines.append("无已通过项。")
            lines.append("")
        else:
            lines.append("| 编号 | 名称 | 类型 | 路径 | 说明 |")
            lines.append("|:--:|:--|:--:|:--|:--|")
            for r in passed_items:
                msg_text = "; ".join(r.messages) if r.messages else "—"
                lines.append(
                    f"| {r.item_id} | {r.name} | {r.item_type} | `{r.path}` | {msg_text} |"
                )
            lines.append("")

        # 下一步行动
        lines.append("## 📋 下一步行动")
        lines.append("")
        if not failed_items:
            lines.append("所有交付物已就位，可以执行 `--pack` 打包提交。")
        else:
            lines.append("按以下顺序处理缺失项：")
            lines.append("")
            # 按优先级排序：error 优先于 warning
            ordered = sorted(failed_items, key=lambda x: 0 if x.severity == "error" else 1)
            for idx, r in enumerate(ordered, 1):
                guidance = self.MISSING_ITEM_GUIDANCE.get(r.item_id, "—")
                lines.append(f"{idx}. **[{r.item_id}] {r.name}** — {guidance}")
            lines.append("")
            lines.append(
                "> 处理完成后重新运行 `python scripts/ci/validate_submission.py --check` 验证。"
            )
        lines.append("")

        output = Path(output_path)
        output.parent.mkdir(parents=True, exist_ok=True)
        output.write_text("\n".join(lines), encoding="utf-8")
        print(f"\n📝 缺失项清单已生成: {output}")


def prepare_submission(manifest_path: str, project_root: str = ".") -> None:
    """准备可自动生成的提交物

    创建 dist/ 目录、生成缺失项报告、输出人工交付物检查清单。

    Args:
        manifest_path: 清单文件路径
        project_root: 项目根目录
    """
    root = Path(project_root)

    # 1. 创建 dist/ 目录
    dist_dir = root / "dist"
    dist_dir.mkdir(exist_ok=True)
    print(f"✅ 确保目录存在: {dist_dir}")

    # 2. 运行校验
    validator = SubmissionValidator(manifest_path, project_root)
    validator.validate_all()

    # 3. 生成缺失项报告
    report_path = root / "results" / "reports" / "submission_validation_report.md"
    validator.generate_report(str(report_path))

    # 4. 输出人工交付物检查清单
    print("\n" + "=" * 60)
    print("📋 人工交付物检查清单")
    print("=" * 60)

    manual_items: list[dict[str, str]] = []
    for r in validator.results:
        if not r.passed and r.item_id in (
            "WHITEPAPER",
            "PRESENTATION",
            "DEMO_VIDEO",
            "CODE_REPO",
            "CODE_ARCHIVE",
            "SUMMARY_REPORT",
            "REGISTRATION_FORM",
        ):
            guidance = SubmissionValidator.MISSING_ITEM_GUIDANCE.get(r.item_id, "—")
            manual_items.append({"id": r.item_id, "name": r.name, "guidance": guidance})

    if manual_items:
        for item in manual_items:
            print(f"\n  [{item['id']}] {item['name']}")
            print(f"    → {item['guidance']}")
    else:
        print("\n  所有交付物已就位！可以执行 --pack 打包。")

    print("\n" + "=" * 60)
    print("📦 自动准备完成：")
    print("  - dist/ 目录已创建")
    print(f"  - 缺失项报告已生成: {report_path}")
    print("  - 人工交付物检查清单已输出（见上方）")
    print()
    print("下一步：")
    print("  1. 完成上方人工交付物")
    print("  2. 运行 --check 验证所有项通过")
    print("  3. 运行 --pack 生成最终压缩包")
    print("=" * 60)


def _is_excluded(rel: str, exclude_list: list[str]) -> bool:
    """判断相对路径是否命中 exclude 列表。

    8.11 修复：除顶层前缀匹配外，支持任意路径段匹配（如 exclude "__pycache__"
    可排除 src/__pycache__/x.pyc 等嵌套路径）；以 "/" 结尾的 exclude 项按目录段匹配。
    """
    parts = rel.split("/")
    for exc in exclude_list:
        exc = exc.rstrip("/")
        if not exc:
            continue
        if rel == exc or rel.startswith(exc + "/"):
            return True
        if exc in parts:  # 嵌套路径段匹配（目录名/文件名）
            return True
    return False


def _add_zip_entry(
    zipf: zipfile.ZipFile, root_path: Path, file: Path, exclude_list: list[str]
) -> None:
    """以根目录相对路径将文件写入 zip，命中 exclude 前缀时跳过"""
    rel = str(file.relative_to(root_path)).replace("\\", "/")
    if not _is_excluded(rel, exclude_list):
        zipf.write(file, rel)


def _build_code_archive(manifest_path: str, project_root: str) -> Path | None:
    """按清单 CODE_ARCHIVE 项的 include/exclude 规则构建代码压缩包。

    仅打包 include 列表中的路径，并排除 exclude 列表中的路径，
    避免整仓递归（含 .git/、dist/）导致的体积失控与自包含问题。

    Args:
        manifest_path: 清单文件路径
        project_root: 项目根目录

    Returns:
        生成的代码压缩包路径；清单中不存在 CODE_ARCHIVE 项时返回 None
    """
    with open(manifest_path, encoding="utf-8") as f:
        manifest = yaml.safe_load(f)
    root = Path(project_root)

    for item in manifest["items"]:
        if item["id"] != "CODE_ARCHIVE":
            continue
        archive_path = root / item["path"]
        reqs = item.get("requirements", {})
        include_list = reqs.get("include", [])
        exclude_list = [exc.rstrip("/") for exc in reqs.get("exclude", [])]

        archive_path.parent.mkdir(parents=True, exist_ok=True)

        with zipfile.ZipFile(archive_path, "w", zipfile.ZIP_DEFLATED) as zipf:
            if include_list:
                for inc in include_list:
                    inc_clean = inc.rstrip("/")
                    inc_path = root / inc_clean
                    if not inc_path.exists():
                        continue
                    if inc_path.is_file():
                        _add_zip_entry(zipf, root, inc_path, exclude_list)
                    else:
                        for file in inc_path.rglob("*"):
                            if file.is_file():
                                _add_zip_entry(zipf, root, file, exclude_list)
            else:
                for file in root.rglob("*"):
                    if file.is_file() and not file.as_posix().startswith("dist/"):
                        _add_zip_entry(zipf, root, file, exclude_list)
        print(f"  📦 代码压缩包: {archive_path}")
        return archive_path
    return None


def package_submission(
    manifest_path: str, project_root: str = ".", skip_items: list[str] | None = None
) -> None:
    """打包提交物

    Args:
        manifest_path: 清单文件路径
        project_root: 项目根目录
        skip_items: 跳过的提交物 id 列表（Issue #860）
    """
    validator = SubmissionValidator(manifest_path, project_root, skip_items=skip_items)

    # 先构建代码压缩包（dist/quantum-rl-scheduler-v9.1.zip），再校验，
    # 解决 CODE_ARCHIVE "先有鸡还是先有蛋" 的循环依赖
    _build_code_archive(manifest_path, project_root)

    if not validator.validate_all():
        print("\n❌ 校验失败，拒绝打包")
        sys.exit(1)

    print("\n📦 开始打包提交物...")

    # 创建输出目录
    output_dir = Path(project_root) / "dist"
    output_dir.mkdir(exist_ok=True)

    # 生成输出文件名
    version = validator.manifest["submission"]["version"]
    date_str = datetime.now().strftime("%Y%m%d")
    output_file = output_dir / f"submission_{version}_{date_str}.zip"
    output_resolved = output_file.resolve()

    # 8.11 修复：外层包也应用 CODE_ARCHIVE 的 exclude 规则，
    # 防止 docs/award_roadmap.md 等内部文档泄漏进外层 zip
    code_archive_excludes: list[str] = []
    for item in validator.manifest["items"]:
        if item["id"] == "CODE_ARCHIVE":
            code_archive_excludes = [
                exc.rstrip("/") for exc in item.get("requirements", {}).get("exclude", [])
            ]
            break

    # 创建 ZIP 文件
    with zipfile.ZipFile(output_file, "w", zipfile.ZIP_DEFLATED) as zipf:
        for item in validator.manifest["items"]:
            # git_tag 项（如 CODE_REPO, path="."）是整仓引用，
            # 打包时跳过，避免递归包含 .git/ 与 dist/ 输出文件本身导致体积失控
            if item["type"] == "git_tag":
                continue
            path = Path(project_root) / item["path"]
            if path.exists():
                if path.is_file():
                    rel = item["path"].replace("\\", "/")
                    if _is_excluded(rel, code_archive_excludes):
                        print(f"  ⏭️ 跳过(内部文档): {item['path']}")
                        continue
                    zipf.write(path, item["path"])
                    print(f"  ✅ 添加: {item['path']}")
                elif path.is_dir():
                    for file in path.rglob("*"):
                        if file.is_file() and file.resolve() != output_resolved:
                            arcname = str(file.relative_to(project_root)).replace("\\", "/")
                            if _is_excluded(arcname, code_archive_excludes):
                                continue
                            zipf.write(file, arcname)
                    print(f"  ✅ 添加目录: {item['path']}")

        # 8.13 round9 审查 P2：外层 zip 无任何指引，评委解压后找不到代码（代码在
        # 嵌套的 dist/quantum-rl-scheduler-v9.1.zip 内）。顶层添加"先读我"说明。
        readme_first = (
            "量子RL驱动的天衍云平台智能调度系统 - 提交包说明\n"
            "==================================================\n"
            "\n"
            "本压缩包为完整参赛提交物，目录结构如下：\n"
            "\n"
            "  docs/                      参赛总结报告 / 参赛报名表 / 技术白皮书(PDF) 等\n"
            "  deliverable_models/        答辩PPT + 权威模型（PPO 16维 / 编译层 / 公平17维 / MAPPO）\n"
            "  dist/quantum-rl-scheduler-v9.1.zip   源代码压缩包（需再解压一层）\n"
            "  results/reports/           实验报告（策略对比 / 消融 / 统计显著性 / 真机验证等）\n"
            "  演示视频_量子RL调度系统.mp4  演示视频\n"
            "\n"
            "快速开始：\n"
            "  1. 解压 dist/quantum-rl-scheduler-v9.1.zip 得到源代码\n"
            "  2. 环境要求：Python 3.10-3.12\n"
            "  3. 安装依赖：pip install -r requirements.lock（推荐）或 pip install -r requirements.txt\n"
            "  4. 运行演示：python scripts/cli.py simulate --episodes 5\n"
            "\n"
            "核心结论（权威口径见源代码内 config/statistics.yaml）：\n"
            "  - PPO 综合调度奖励 vs 真实 FCFS 提升 +20.2%（N=250, Welch t p=7.56e-12）\n"
            "  - 等待时间 -14.0%（N=250 配对检验 p=8.9e-06）\n"
            "  - 真机 SDK 调用 315 次 100% 成功（可用性验证）\n"
            "\n"
            "更多信息见源代码内 README.md。\n"
        )
        zipf.writestr("README_先读我.txt", readme_first)
        print("  ✅ 添加: README_先读我.txt（外层指引）")

    print(f"\n✅ 打包完成: {output_file}")
    print(f"   文件大小: {output_file.stat().st_size / (1024 * 1024):.1f}MB")
    return output_file


def _rename_final_package(
    result_path: Path, final_name: str, final_dir: str, project_root: str
) -> None:
    """按比赛规范命名最终压缩包（提报单位-选题名称-作品名称）。"""
    root = Path(project_root)
    out_dir = Path(final_dir) if final_dir else root
    out_dir.mkdir(parents=True, exist_ok=True)
    safe_name = "".join(c for c in final_name if c not in '<>:"/\\|?*')
    target = out_dir / f"{safe_name}.zip"
    import shutil

    shutil.copy2(result_path, target)
    print(f"\n✅ 最终提交包（比赛命名）: {target}")
    print(f"   大小: {target.stat().st_size / (1024 * 1024):.1f}MB")
    print("   命名规范: 提报单位（学校全称）－选题名称－作品名称")


def main() -> None:
    """主函数"""
    parser = argparse.ArgumentParser(
        description="M5 最终提交物校验与打包工具",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  # 仅校验
  python scripts/ci/validate_submission.py --check

  # 校验并打包
  python scripts/ci/validate_submission.py --pack

  # 校验并生成缺失项清单报告（Issue #168）
  python scripts/ci/validate_submission.py --check --report results/reports/submission_validation_report.md

  # 自定义路径
  python scripts/ci/validate_submission.py --check --manifest config/submission_manifest.yaml --project-root .
        """,
    )

    parser.add_argument("--check", action="store_true", help="仅校验提交物")
    parser.add_argument("--pack", action="store_true", help="校验并打包提交物")
    parser.add_argument(
        "--prepare",
        action="store_true",
        help="准备可自动生成的提交物（dist 目录、缺失项报告、检查清单）",
    )
    parser.add_argument(
        "--report",
        type=str,
        default=None,
        help="将校验结果输出为 Markdown 格式的缺失项清单报告（推荐: results/reports/submission_validation_report.md）",
    )
    parser.add_argument(
        "--manifest",
        type=str,
        default="config/submission_manifest.yaml",
        help="清单文件路径 (默认: config/submission_manifest.yaml)",
    )
    parser.add_argument(
        "--skip-items",
        type=str,
        default="",
        help="逗号分隔的提交物 id 列表，跳过其校验（如 CODE_ARCHIVE,DEMO_VIDEO；"
        "用于 release CI 中跳过仅本地生成的交付物）",
    )
    parser.add_argument(
        "--project-root",
        type=str,
        default=".",
        help="项目根目录 (默认: 当前目录)",
    )
    parser.add_argument(
        "--rename-final",
        type=str,
        default="",
        help="按比赛规范重命名最终压缩包: '提报单位（学校全称）-选题名称-作品名称'，"
        "如 'XX大学-量子+AI双向赋能-量子RL调度系统'。"
        "--pack 完成后将 dist/submission_*.zip 复制为 <值>.zip（顶层目录）。"
        "可选后缀: --final-dir 指定输出目录（默认项目根）。",
    )
    parser.add_argument(
        "--final-dir",
        type=str,
        default="",
        help="--rename-final 输出目录（默认项目根目录）",
    )

    args = parser.parse_args()

    if not args.check and not args.pack and not args.prepare:
        parser.error("必须指定 --check、--pack 或 --prepare 之一")

    skip_ids = [s.strip() for s in args.skip_items.split(",") if s.strip()]

    if args.prepare:
        prepare_submission(args.manifest, args.project_root)
    elif args.pack:
        result_path = package_submission(args.manifest, args.project_root, skip_items=skip_ids)
        if args.rename_final:
            _rename_final_package(result_path, args.rename_final, args.final_dir, args.project_root)
    else:
        validator = SubmissionValidator(args.manifest, args.project_root, skip_items=skip_ids)
        success = validator.validate_all()
        if args.report:
            validator.generate_report(args.report)
        sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()
